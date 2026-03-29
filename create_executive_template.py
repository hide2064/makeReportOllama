"""
create_executive_template.py
コンサルファーム最上位グレードの PPTX テンプレートを生成する。

生成ファイル:
  data/template_executive.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

DATA_DIR = Path(__file__).parent / "data"
DATA_DIR.mkdir(exist_ok=True)

# ── スライドサイズ (16:9 ワイド) ─────────────────────────────────
W = 12_192_000
H =  6_858_000

# ── プレミアムカラーパレット ─────────────────────────────────────
C_INK       = RGBColor(0x0A, 0x12, 0x28)   # 最深ネイビー (背景・ヘッダー)
C_NAVY      = RGBColor(0x0F, 0x20, 0x44)   # ダークネイビー
C_NAVY_MID  = RGBColor(0x16, 0x32, 0x60)   # ミッドネイビー (サイドバー)
C_NAVY_LITE = RGBColor(0x1E, 0x45, 0x7A)   # ライトネイビー (テーブルヘッダー)
C_GOLD      = RGBColor(0xD4, 0x94, 0x1A)   # ゴールド (主アクセント)
C_GOLD_PALE = RGBColor(0xF5, 0xD9, 0x8A)   # ペールゴールド (ハイライト)
C_TEAL      = RGBColor(0x00, 0xA3, 0x9A)   # ティール (サブアクセント)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_OFFWHITE  = RGBColor(0xF6, 0xF8, 0xFC)   # 本文背景
C_SURFACE   = RGBColor(0xEE, 0xF3, 0xFA)   # カード背景
C_TEXT      = RGBColor(0x12, 0x1A, 0x2E)   # 本文テキスト
C_TEXT_MID  = RGBColor(0x3A, 0x4E, 0x68)   # セカンダリテキスト
C_TEXT_LITE = RGBColor(0x6B, 0x7E, 0x96)   # ミュートテキスト
C_RULE      = RGBColor(0xCF, 0xD9, 0xE8)   # 区切り線
C_RED_SOFT  = RGBColor(0xC0, 0x35, 0x2B)   # 警告/強調


def rect(slide, l, t, w, h, color, line_color=None):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Emu(6_000)
    else:
        s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, size, bold=False, italic=False,
       color=C_TEXT, align=PP_ALIGN.LEFT, font="Arial", wrap=True):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return box


def tb_multi(slide, l, t, w, h, lines, size, bold=False,
             color=C_TEXT, align=PP_ALIGN.LEFT, line_spacing=1.3):
    """複数行テキストボックス。lines = [(text, bold, color)] のリスト。"""
    from pptx.oxml.ns import qn
    from lxml import etree
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (text, line_bold, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = line_bold
        run.font.color.rgb = line_color
        run.font.name  = "Arial"
    return box


# ────────────────────────────────────────────────────────────────
def build_template():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    # ── レイアウト取得 ──────────────────────────────────────────
    blank = next(
        (l for l in prs.slide_layouts if not l.placeholders),
        prs.slide_layouts[-1]
    )

    # ════════════════════════════════════════════════════════════
    # SLIDE 1 — カバースライド (インパクトデザイン)
    # ════════════════════════════════════════════════════════════
    s0 = prs.slides.add_slide(blank)
    for ph in list(s0.placeholders):
        ph._element.getparent().remove(ph._element)

    # ── フルブリード背景 ─────────────────────────────────────
    rect(s0, 0, 0, W, H, C_INK)

    # ── 斜めゴールドストライプ（右上→左下へ走る装飾帯） ──────
    # 幅広の斜め帯は四角形+回転で近似
    STRIPE_W = 1_900_000
    stripe = s0.shapes.add_shape(1, Emu(W - STRIPE_W - 500_000), Emu(-200_000),
                                  Emu(STRIPE_W), Emu(H + 400_000))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = C_GOLD
    stripe.line.fill.background()
    stripe.rotation = -8.0

    # さらに細い半透明帯 (ティール) を重ねる
    stripe2 = s0.shapes.add_shape(1, Emu(W - STRIPE_W - 1_400_000), Emu(-200_000),
                                   Emu(400_000), Emu(H + 400_000))
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = C_TEAL
    stripe2.line.fill.background()
    stripe2.rotation = -8.0

    # ── 左辺の縦アクセントライン ─────────────────────────────
    rect(s0, 0, 0, 32_000, H, C_GOLD)

    # ── 上部ゴールドルール ────────────────────────────────────
    rect(s0, 32_000, 1_600_000, W - 32_000, 14_000, RGBColor(0xD4, 0x94, 0x1A))

    # ── ロゴ/会社名エリア ──────────────────────────────────────
    tb(s0, 180_000, 680_000, 3_000_000, 600_000,
       "SALES INTELLIGENCE", 9, bold=True,
       color=C_GOLD, font="Arial")
    rect(s0, 180_000, 1_200_000, 600_000, 8_000, C_GOLD)

    # ── メインタイトル ────────────────────────────────────────
    tb(s0, 180_000, 1_700_000, W - 2_200_000, 1_200_000,
       "{{report_title}}", 30, bold=True,
       color=C_WHITE, font="Arial")

    # ── サブタイトルライン ────────────────────────────────────
    tb(s0, 180_000, 2_980_000, 4_000_000, 600_000,
       "月次経営報告書  ―  AI 生成", 13,
       color=C_GOLD_PALE, font="Arial")

    # ── 作成日エリア ──────────────────────────────────────────
    tb(s0, 180_000, 3_700_000, 3_000_000, 380_000,
       "{{report_date}}", 10,
       color=C_TEXT_LITE, font="Arial")

    # ── 下部区切り＋キャッチコピー ───────────────────────────
    rect(s0, 0, H - 520_000, W, 520_000, C_NAVY)
    rect(s0, 0, H - 520_000, W, 10_000, C_GOLD)
    tb(s0, 180_000, H - 450_000, W - 360_000, 380_000,
       "CONFIDENTIAL   |   本資料は機密情報を含みます。取り扱いにご注意ください。",
       8, color=C_TEXT_LITE, font="Arial")

    # ════════════════════════════════════════════════════════════
    # SLIDE 2 — エグゼクティブサマリー (メインコンテンツ)
    # ════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    for ph in list(s1.placeholders):
        ph._element.getparent().remove(ph._element)

    # ── 背景 ─────────────────────────────────────────────────
    rect(s1, 0, 0, W, H, C_OFFWHITE)

    # ── トップヘッダーバー ────────────────────────────────────
    HDR_H = 760_000
    rect(s1, 0, 0, W, HDR_H, C_INK)

    # ヘッダー内: ゴールド縦ライン
    rect(s1, 0, 0, 22_000, HDR_H, C_GOLD)

    # ヘッダー内: タイトル
    tb(s1, 80_000, 120_000, W - 2_600_000, HDR_H - 160_000,
       "{{report_title}}", 18, bold=True, color=C_WHITE)

    # ヘッダー内: 右側 BADGE (MONTHLY REPORT)
    badge_l = W - 2_300_000
    rect(s1, badge_l, 180_000, 680_000, 380_000, C_TEAL)
    tb(s1, badge_l, 220_000, 680_000, 340_000,
       "MONTHLY", 7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ヘッダー内: 日付
    tb(s1, badge_l + 720_000, 200_000, 1_400_000, 380_000,
       "{{report_date}}", 9, color=RGBColor(0xAA, 0xBB, 0xCC))

    # ── ヘッダー下のゴールドアクセントライン ──────────────────
    rect(s1, 0, HDR_H, W, 16_000, C_GOLD)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 2カラムコンテンツエリア
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    MARGIN    = 200_000
    GAP       = 160_000
    COL_TOP   = HDR_H + 16_000 + 80_000
    COL_H     = H - COL_TOP - 440_000
    COL_W_L   = int((W - MARGIN * 2 - GAP) * 0.54)   # 左54%
    COL_W_R   = W - MARGIN * 2 - GAP - COL_W_L        # 右46%
    COL_L_L   = MARGIN
    COL_R_L   = MARGIN + COL_W_L + GAP

    # ── 左カラム: サマリー ────────────────────────────────────
    # カードシャドウ代わりに微かに暗い底面
    rect(s1, COL_L_L + 8_000, COL_TOP + 8_000, COL_W_L, COL_H,
         RGBColor(0xD4, 0xDC, 0xEC))
    # カード本体
    rect(s1, COL_L_L, COL_TOP, COL_W_L, COL_H, C_WHITE)

    # 左カラム: トップラベルバー (ネイビー)
    LBL_H = 420_000
    rect(s1, COL_L_L, COL_TOP, COL_W_L, LBL_H, C_NAVY)
    # ラベルバー左辺の太いゴールドライン
    rect(s1, COL_L_L, COL_TOP, 18_000, LBL_H, C_GOLD)
    # ラベルテキスト
    tb(s1, COL_L_L + 60_000, COL_TOP + 80_000,
       COL_W_L - 100_000, LBL_H - 100_000,
       "EXECUTIVE SUMMARY  /  売上サマリー",
       9, bold=True, color=C_WHITE)

    # 左カラム: セクション境界ライン
    rect(s1, COL_L_L + 60_000, COL_TOP + LBL_H + 100_000,
         280_000, 8_000, C_GOLD)

    # 左カラム: 本文プレースホルダー
    tb(s1, COL_L_L + 60_000, COL_TOP + LBL_H + 180_000,
       COL_W_L - 120_000, COL_H - LBL_H - 240_000,
       "{{summary_text}}", 10.5, color=C_TEXT)

    # ── 右カラム: 課題・分析 ──────────────────────────────────
    # カード
    rect(s1, COL_R_L + 8_000, COL_TOP + 8_000, COL_W_R, COL_H,
         RGBColor(0xD4, 0xDC, 0xEC))
    rect(s1, COL_R_L, COL_TOP, COL_W_R, COL_H, C_SURFACE)

    # 右カラム: トップラベルバー (ゴールド)
    rect(s1, COL_R_L, COL_TOP, COL_W_R, LBL_H, C_GOLD)
    rect(s1, COL_R_L, COL_TOP, 18_000, LBL_H, C_NAVY)
    tb(s1, COL_R_L + 60_000, COL_TOP + 80_000,
       COL_W_R - 100_000, LBL_H - 100_000,
       "KEY INSIGHTS  /  課題・改善策",
       9, bold=True, color=C_INK)

    # 右カラム: セクション境界ライン
    rect(s1, COL_R_L + 60_000, COL_TOP + LBL_H + 100_000,
         280_000, 8_000, C_NAVY)

    # 右カラム: 本文プレースホルダー
    tb(s1, COL_R_L + 60_000, COL_TOP + LBL_H + 180_000,
       COL_W_R - 120_000, COL_H - LBL_H - 240_000,
       "{{analysis_text}}", 10.5, color=C_TEXT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # フッター
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    FTR_T = H - 400_000
    rect(s1, 0, FTR_T, W, 400_000, C_INK)
    rect(s1, 0, FTR_T, W, 10_000, C_GOLD)

    # フッター左: システム名
    tb(s1, 200_000, FTR_T + 90_000, W // 2 - 200_000, 260_000,
       "Sales Report  |  Powered by Ollama (Local LLM)  ·  AI Generated",
       8, color=C_TEXT_LITE)

    # フッター右: 機密表示
    tb(s1, W // 2, FTR_T + 90_000, W // 2 - 200_000, 260_000,
       "STRICTLY CONFIDENTIAL  —  社外秘",
       8, bold=True, color=C_GOLD_PALE, align=PP_ALIGN.RIGHT)

    # ページ番号ドット
    rect(s1, W - 240_000, FTR_T + 140_000, 80_000, 80_000, C_GOLD)

    # ════════════════════════════════════════════════════════════
    # SLIDE 3 — アペンディックス用白紙スライド (追加スライドの前置き)
    # ════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    for ph in list(s2.placeholders):
        ph._element.getparent().remove(ph._element)

    rect(s2, 0, 0, W, H, C_INK)
    rect(s2, 0, 0, 22_000, H, C_GOLD)
    rect(s2, 0, H - 520_000, W, 520_000, C_NAVY)
    rect(s2, 0, H - 520_000, W, 10_000, C_GOLD)

    # 中央テキスト
    tb(s2, 0, H // 2 - 600_000, W, 500_000,
       "APPENDIX", 48, bold=True,
       color=C_GOLD, align=PP_ALIGN.CENTER, font="Arial")
    tb(s2, 0, H // 2 - 50_000, W, 400_000,
       "詳細データ・グラフ", 16,
       color=C_TEXT_LITE, align=PP_ALIGN.CENTER)

    rect(s2, W // 2 - 600_000, H // 2 + 420_000, 1_200_000, 10_000, C_GOLD)

    tb(s2, 200_000, H - 440_000, W - 400_000, 280_000,
       "STRICTLY CONFIDENTIAL  —  社外秘", 8,
       bold=True, color=C_GOLD_PALE, align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════════════════════
    # 保存
    # ════════════════════════════════════════════════════════════
    out = DATA_DIR / "template_executive.pptx"
    prs.save(out)
    print(f"✓  エグゼクティブテンプレート生成完了: {out}")
    print(f"   スライド数: {len(prs.slides)} 枚")
    print(f"     1. カバー（インパクトデザイン）")
    print(f"     2. エグゼクティブサマリー（2カラム）")
    print(f"     3. アペンディックス区切り")


if __name__ == "__main__":
    build_template()
