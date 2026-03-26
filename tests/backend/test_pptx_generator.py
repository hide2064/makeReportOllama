"""
test_pptx_generator.py — pptx_generator サービスのユニットテスト
"""
import io
import os

import pandas as pd
import pytest
from pptx import Presentation

from services.pptx_generator import (
    _add_blank_slide,
    _add_chart_slide,
    _add_table_slide,
    generate_pptx,
)

# プロジェクトルートからテンプレートパスを解決
TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__), "../../data/template.pptx"
)


@pytest.fixture()
def output_path(tmp_path):
    return str(tmp_path / "output.pptx")


def test_generates_file(output_path):
    generate_pptx(
        template_path=TEMPLATE_PATH,
        output_path=output_path,
        summary_text="テストサマリー",
        analysis_text="テスト分析",
        period="2025/01/01 ～ 2025/03/31",
    )
    assert os.path.exists(output_path)


def test_placeholders_replaced(output_path):
    generate_pptx(
        template_path=TEMPLATE_PATH,
        output_path=output_path,
        summary_text="サマリー本文テキスト",
        analysis_text="分析本文テキスト",
        period="2025/01/01 ～ 2025/03/31",
    )
    prs = Presentation(output_path)
    all_text = " ".join(
        run.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    assert "{{summary_text}}"  not in all_text
    assert "{{analysis_text}}" not in all_text
    assert "{{report_title}}"  not in all_text
    assert "サマリー本文テキスト"  in all_text
    assert "分析本文テキスト"    in all_text


def test_template_not_found(output_path):
    with pytest.raises(FileNotFoundError):
        generate_pptx(
            template_path="/nonexistent/template.pptx",
            output_path=output_path,
            summary_text="",
            analysis_text="",
            period="",
        )


# ── _add_table_slide ────────────────────────────────────────────

def _make_pivot():
    """テスト用の四半期×商品クロス集計 DataFrame を返す。"""
    data = {
        "2024Q1": [100_000, 80_000],
        "2024Q2": [120_000, 90_000],
    }
    return pd.DataFrame(data, index=["商品A", "商品B"])


def test_add_table_slide_appends_slide():
    prs = Presentation(TEMPLATE_PATH)
    slide_count_before = len(prs.slides)
    _add_table_slide(prs, _make_pivot(), "商品別売上表", "商品名")
    assert len(prs.slides) == slide_count_before + 1


def test_add_table_slide_skips_on_empty():
    prs = Presentation(TEMPLATE_PATH)
    slide_count_before = len(prs.slides)
    empty_df = pd.DataFrame()
    _add_table_slide(prs, empty_df, "商品別売上表", "商品名")
    assert len(prs.slides) == slide_count_before  # スキップされること


def test_add_table_slide_cell_values():
    prs = Presentation(TEMPLATE_PATH)
    _add_table_slide(prs, _make_pivot(), "商品別売上表", "商品名")
    slide = prs.slides[-1]
    # テーブルシェイプを探す
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    tbl = tables[0].table
    # ヘッダー行 (0行目) の最初のセルが row_header であること
    assert tbl.cell(0, 0).text_frame.paragraphs[0].runs[0].text == "商品名"
    # データ行に商品名が含まれること
    row_labels = [tbl.cell(i, 0).text_frame.paragraphs[0].runs[0].text for i in range(1, 3)]
    assert "商品A" in row_labels
    assert "商品B" in row_labels


def test_add_table_slide_generic_region(tmp_path):
    """地域別ピボットでも同じ関数が使えること。"""
    prs = Presentation(TEMPLATE_PATH)
    region_pivot = pd.DataFrame(
        {"2024Q1": [200_000, 150_000]}, index=["東京", "大阪"]
    )
    _add_table_slide(prs, region_pivot, "地域別売上表", "地域")
    slide = prs.slides[-1]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    tbl = tables[0].table
    assert tbl.cell(0, 0).text_frame.paragraphs[0].runs[0].text == "地域"


# ── _add_chart_slide ────────────────────────────────────────────

def _make_chart_data():
    monthly = {"2025-01": 1_000_000, "2025-02": 1_200_000, "2025-03": 900_000}
    products = {"商品A": 1_500_000, "商品B": 900_000, "商品C": 700_000}
    return monthly, products


def test_add_chart_slide_bar_appends_slide():
    prs = Presentation(TEMPLATE_PATH)
    monthly, products = _make_chart_data()
    slide_count_before = len(prs.slides)
    _add_chart_slide(prs, monthly, products, product_chart_type="bar")
    assert len(prs.slides) == slide_count_before + 1


def test_add_chart_slide_pie_appends_slide():
    prs = Presentation(TEMPLATE_PATH)
    monthly, products = _make_chart_data()
    slide_count_before = len(prs.slides)
    _add_chart_slide(prs, monthly, products, product_chart_type="pie")
    assert len(prs.slides) == slide_count_before + 1


def test_add_chart_slide_skips_on_empty():
    prs = Presentation(TEMPLATE_PATH)
    slide_count_before = len(prs.slides)
    _add_chart_slide(prs, {}, {})
    assert len(prs.slides) == slide_count_before


def test_add_chart_slide_with_margin():
    prs = Presentation(TEMPLATE_PATH)
    monthly, products = _make_chart_data()
    margin = {"2025-01": 45.2, "2025-02": 48.0, "2025-03": 42.5}
    slide_count_before = len(prs.slides)
    _add_chart_slide(prs, monthly, products, monthly_margin=margin, product_chart_type="bar")
    assert len(prs.slides) == slide_count_before + 1


# ── slide_options ───────────────────────────────────────────────

def test_generate_pptx_slide_options_all_off(output_path):
    """slide_options でテーブル・グラフをすべてオフにするとテンプレートのスライド数のまま。"""
    prs_template = Presentation(TEMPLATE_PATH)
    n_template = len(prs_template.slides)

    generate_pptx(
        template_path=TEMPLATE_PATH,
        output_path=output_path,
        summary_text="サマリー",
        analysis_text="分析",
        period="2025/01 ～ 2025/03",
        slide_options={
            "product_table": False,
            "region_table": False,
            "rep_table": False,
            "chart": False,
        },
    )
    prs_out = Presentation(output_path)
    assert len(prs_out.slides) == n_template


def test_generate_pptx_slide_options_region_rep(output_path):
    """region_table と rep_table をオンにするとテンプレート + 2 スライド追加される。"""
    prs_template = Presentation(TEMPLATE_PATH)
    n_template = len(prs_template.slides)

    region_pivot = pd.DataFrame({"2024Q1": [100_000]}, index=["東京"])
    rep_pivot    = pd.DataFrame({"2024Q1": [100_000]}, index=["田中"])

    generate_pptx(
        template_path=TEMPLATE_PATH,
        output_path=output_path,
        summary_text="サマリー",
        analysis_text="分析",
        period="2025/01 ～ 2025/03",
        quarterly_region_pivot=region_pivot,
        quarterly_rep_pivot=rep_pivot,
        slide_options={
            "product_table": False,
            "region_table": True,
            "rep_table": True,
            "chart": False,
        },
    )
    prs_out = Presentation(output_path)
    assert len(prs_out.slides) == n_template + 2
