"""
create_sample_report.py
RAG 登録用サンプル売上報告資料 (5スライド) を生成する。
  Slide 1: タイトルページ
  Slide 2: 売上実績表
  Slide 3: 売上分析（トレンド・商品・地域）
  Slide 4: 現状分析 — メリット・デメリット・今後の方針
  Slide 5: まとめ
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

OUT = Path(__file__).parent / "data" / "sample_report_2024Q1.pptx"

# ── カラー ────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1B, 0x2E, 0x4C)
C_NAVY2  = RGBColor(0x2C, 0x4A, 0x7A)
C_GOLD   = RGBColor(0xC4, 0x97, 0x3E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY   = RGBColor(0x55, 0x65, 0x7A)
C_LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
C_LGOLD  = RGBColor(0xFD, 0xF6, 0xEC)
C_GREEN  = RGBColor(0x16, 0x75, 0x3A)
C_RED    = RGBColor(0xCC, 0x28, 0x28)
C_BLUE   = RGBColor(0x1A, 0x56, 0xA0)

W = 12_192_000
H =  6_858_000


def rgb(s):
    s.fill.solid()
    return s.fill.fore_color


def rect(slide, l, t, w, h, color, no_line=True):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if no_line: s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, size, bold=False, color=C_TEXT,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Arial"
    return box


def add_para(tf, text, size, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as Pts
    p = tf.add_paragraph(); p.alignment = align
    if space_before:
        p.space_before = Pts(space_before)
    r = p.add_run(); r.text = text
    r.font.size = Pts(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Arial"
    return p


def header_band(slide, title_text, subtitle=""):
    rect(slide, 0, 0, W, 680_000, C_NAVY)
    rect(slide, 0, 0, 160_000, 680_000, C_GOLD)
    tb(slide, 220_000, 100_000, W - 500_000, 480_000,
       title_text, 18, bold=True, color=C_WHITE)
    if subtitle:
        tb(slide, W - 2_200_000, 100_000, 2_100_000, 480_000,
           subtitle, 9, color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.RIGHT)


def footer_band(slide, note="社外秘 — 取扱注意"):
    rect(slide, 0, H - 360_000, W, 360_000, C_NAVY)
    rect(slide, 0, H - 360_000, W, 26_000, C_GOLD)
    tb(slide, 360_000, H - 310_000, W - 720_000, 280_000,
       note, 8, color=RGBColor(0xAA, 0xB8, 0xCC))


# ────────────────────────────────────────────────────────────
#  Slide 1: タイトル
# ────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, C_NAVY)
    # ゴールド縦帯
    rect(slide, 0, 0, 220_000, H, C_GOLD)
    # 装飾横線
    rect(slide, 220_000, H // 2 - 10_000, W - 220_000, 20_000, C_GOLD)

    tb(slide, 400_000, 1_200_000, W - 500_000, 1_000_000,
       "2024年 第1四半期", 16, color=RGBColor(0xCC, 0xD6, 0xE8))
    tb(slide, 400_000, 2_000_000, W - 500_000, 1_200_000,
       "売上報告書", 42, bold=True, color=C_WHITE)
    tb(slide, 400_000, 3_200_000, W - 500_000, 600_000,
       "2024年1月 〜 3月 実績", 16, color=RGBColor(0xCC, 0xD6, 0xE8))

    tb(slide, 400_000, H - 1_400_000, W - 500_000, 400_000,
       "株式会社サンプル商事　営業企画部", 11,
       color=RGBColor(0xAA, 0xB8, 0xCC))
    tb(slide, 400_000, H - 1_000_000, W - 500_000, 400_000,
       "作成日: 2024年4月5日", 10,
       color=RGBColor(0xAA, 0xB8, 0xCC))
    footer_band(slide, "Confidential  |  Sales Report Q1 2024")


# ────────────────────────────────────────────────────────────
#  Slide 2: 売上実績表
# ────────────────────────────────────────────────────────────
def slide_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, RGBColor(0xF7, 0xF8, 0xFA))
    header_band(slide, "売上実績（2024年Q1）", "単位: 万円")

    # ── サマリーボックス ────────────────────────────────────
    BOX_T = 780_000; BOX_H = 700_000
    boxes = [
        ("総売上",      "3,847万円",  "+12.3%", C_GREEN),
        ("総販売数量",  "1,234 個",   "+8.1%",  C_GREEN),
        ("平均利益率",  "52.4%",      "+1.2pt", C_GREEN),
        ("最高売上商品","プレミアムP","¥1,520万", C_BLUE),
    ]
    bw = (W - 800_000) // 4
    for i, (lbl, val, sub, sc) in enumerate(boxes):
        bx = 400_000 + i * (bw + 80_000)
        rect(slide, bx, BOX_T, bw, BOX_H, C_WHITE)
        rect(slide, bx, BOX_T, bw, 12_000, sc)
        tb(slide, bx + 80_000, BOX_T + 60_000, bw - 100_000, 260_000,
           lbl, 9, color=C_GRAY)
        tb(slide, bx + 80_000, BOX_T + 280_000, bw - 100_000, 320_000,
           val, 15, bold=True, color=C_TEXT)
        tb(slide, bx + 80_000, BOX_T + 560_000, bw - 100_000, 180_000,
           sub, 9, bold=True, color=sc)

    # ── 月別明細テーブル ────────────────────────────────────
    TBL_T = 1_620_000
    headers = ["月", "プレミアムP", "スタンダードP", "コンサルティング", "保守S", "合計", "利益率"]
    rows = [
        ["1月", "520万", "310万", "480万", "95万", "1,405万", "54.2%"],
        ["2月", "480万", "290万", "420万", "88万", "1,278万", "51.8%"],
        ["3月", "520万", "330万", "220万", "94万", "1,164万", "51.1%"],
        ["合計", "1,520万", "930万", "1,120万", "277万", "3,847万", "52.4%"],
    ]
    col_w = [320_000, 1_380_000, 1_380_000, 1_700_000, 900_000, 1_150_000, 900_000]
    row_h = 440_000
    total_w = sum(col_w)

    for ci, (hdr, cw) in enumerate(zip(headers, col_w)):
        cx = 400_000 + sum(col_w[:ci])
        rect(slide, cx, TBL_T, cw, row_h, C_NAVY2)
        tb(slide, cx + 30_000, TBL_T + 80_000, cw - 40_000, row_h - 80_000,
           hdr, 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        ry = TBL_T + row_h * (ri + 1)
        bg = C_LGRAY if ri % 2 == 0 else C_WHITE
        if ri == 3:
            bg = RGBColor(0xE8, 0xEE, 0xF4)
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            cx = 400_000 + sum(col_w[:ci])
            rect(slide, cx, ry, cw, row_h, bg)
            fc = C_GREEN if (ri == 3 and ci == 5) else (
                 C_RED if cell.startswith("-") else C_TEXT)
            tb(slide, cx + 30_000, ry + 80_000, cw - 40_000, row_h - 80_000,
               cell, 9.5, bold=(ri == 3), color=fc, align=PP_ALIGN.CENTER)

    footer_band(slide)


# ────────────────────────────────────────────────────────────
#  Slide 3: 売上分析
# ────────────────────────────────────────────────────────────
def slide_analysis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, RGBColor(0xF7, 0xF8, 0xFA))
    header_band(slide, "売上分析", "2024年Q1")

    # ── 左列: 商品別構成 ────────────────────────────────────
    COL_T = 820_000; COL_H = H - COL_T - 400_000
    LW = (W - 800_000) // 2 - 60_000
    RX = 400_000 + LW + 120_000

    rect(slide, 400_000, COL_T, LW, COL_H, C_WHITE)
    rect(slide, 400_000, COL_T, LW, 380_000, C_NAVY2)
    tb(slide, 520_000, COL_T + 80_000, LW, 300_000,
       "● 商品別売上構成", 11, bold=True, color=C_WHITE)

    products = [
        ("プレミアムプラン",    "1,520万", 39.5, C_BLUE),
        ("コンサルティング",    "1,120万", 29.1, RGBColor(0x8B, 0x5C, 0xF6)),
        ("スタンダードプラン",    "930万", 24.2, RGBColor(0x0E, 0x9F, 0x6E)),
        ("保守サポート",          "277万",  7.2, C_GOLD),
    ]
    bar_max = 1_600_000
    BY = COL_T + 420_000
    for name, amt, pct, color in products:
        rect(slide, 520_000, BY, int(LW * 0.85 * pct / 40), 160_000, color)
        tb(slide, 520_000, BY + 170_000, LW - 80_000, 200_000,
           f"{name}  {amt} ({pct}%)", 8.5, color=C_TEXT)
        BY += 380_000

    # ── 右列: 地域別 ────────────────────────────────────────
    rect(slide, RX, COL_T, LW, COL_H, C_WHITE)
    rect(slide, RX, COL_T, LW, 380_000, C_GOLD)
    tb(slide, RX + 120_000, COL_T + 80_000, LW, 300_000,
       "● 地域別売上 & 前Q比", 11, bold=True, color=C_NAVY)

    regions = [
        ("東京",   "1,460万", "+15.2%", C_GREEN),
        ("大阪",     "962万",  "+8.7%", C_GREEN),
        ("名古屋",   "693万",  "+5.1%", C_GREEN),
        ("福岡",     "462万",  "+2.3%", C_GREEN),
        ("札幌",     "270万",  "-3.1%", C_RED),
    ]
    RY = COL_T + 420_000
    for region, amt, chg, cc in regions:
        rect(slide, RX + 120_000, RY + 20_000, LW - 200_000, 200_000, C_LGRAY)
        tb(slide, RX + 160_000, RY + 40_000, 800_000, 180_000,
           region, 10, bold=True, color=C_TEXT)
        tb(slide, RX + 900_000, RY + 40_000, 700_000, 180_000,
           amt, 10, color=C_TEXT, align=PP_ALIGN.RIGHT)
        tb(slide, RX + 1_620_000, RY + 40_000, 600_000, 180_000,
           chg, 10, bold=True, color=cc, align=PP_ALIGN.RIGHT)
        RY += 280_000

    footer_band(slide)


# ────────────────────────────────────────────────────────────
#  Slide 4: メリット・デメリット・方針
# ────────────────────────────────────────────────────────────
def slide_swot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, RGBColor(0xF7, 0xF8, 0xFA))
    header_band(slide, "現状分析：強み・課題と今後の方針")

    CARD_T = 800_000; CARD_H = H - CARD_T - 500_000
    CW = (W - 1_000_000) // 3

    cards = [
        ("強み・メリット", C_GREEN, [
            "プレミアムP が全体の 40% を占め高収益を牽引",
            "コンサルTINGの利益率が 65% と最高水準",
            "東京・大阪の主要 2 拠点が安定した売上基盤を形成",
            "前四半期比 +12.3% と 3 期連続の増収を達成",
        ]),
        ("課題・デメリット", C_RED, [
            "札幌エリアが唯一マイナス成長 (▲3.1%)",
            "スタンダードPの利益率が 41% と全商品中最低",
            "コンサルティング単価の季節波動が大きく収益が不安定",
            "担当者 3 名が売上の 55% を占め属人化リスクあり",
        ]),
        ("今後の方針", C_BLUE, [
            "札幌に地域専任担当を配置し Q2 に 10% 回復を目標",
            "スタンダードPの原価見直し / バンドル化で利益率 5pt 改善",
            "コンサルティングの年間契約化を推進し収益を平準化",
            "ナレッジ共有強化と後継担当者育成で属人化リスクを低減",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(cards):
        cx = 400_000 + i * (CW + 100_000)
        rect(slide, cx, CARD_T, CW, CARD_H, C_WHITE)
        rect(slide, cx, CARD_T, CW, 400_000, color)
        # アイコン文字
        icons = ["◎", "▲", "→"]
        tb(slide, cx + 80_000, CARD_T + 60_000, CW - 100_000, 300_000,
           f"{icons[i]}  {title}", 11, bold=True, color=C_WHITE)

        box = slide.shapes.add_textbox(
            Emu(cx + 80_000), Emu(CARD_T + 460_000),
            Emu(CW - 160_000), Emu(CARD_H - 500_000)
        )
        tf = box.text_frame; tf.word_wrap = True
        for j, b in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(6)
            r = p.add_run(); r.text = f"• {b}"
            r.font.size = Pt(9.5); r.font.color.rgb = C_TEXT
            r.font.name = "Arial"

    footer_band(slide)


# ────────────────────────────────────────────────────────────
#  Slide 5: まとめ
# ────────────────────────────────────────────────────────────
def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, C_NAVY)
    rect(slide, 0, 0, 220_000, H, C_GOLD)
    rect(slide, 220_000, H // 2 - 10_000, W - 220_000, 20_000, C_GOLD)

    tb(slide, 400_000, 900_000, W - 500_000, 600_000,
       "まとめ", 28, bold=True, color=C_GOLD)

    points = [
        "Q1 総売上 3,847万円 (前Q比 +12.3%) — 3期連続増収を達成",
        "プレミアムP・コンサルティングの高付加価値商材が収益を牽引",
        "札幌エリアのみ減少 → Q2 に専任担当配置で早期回復を目指す",
        "スタンダードP のコスト構造改善・バンドル化で利益率向上へ",
        "年間契約化・属人化解消を通じて持続可能な収益基盤を構築する",
    ]
    box = slide.shapes.add_textbox(
        Emu(400_000), Emu(1_600_000), Emu(W - 600_000), Emu(3_600_000)
    )
    tf = box.text_frame; tf.word_wrap = True
    for j, pt in enumerate(points):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run(); r.text = f"  ✓  {pt}"
        r.font.size = Pt(12); r.font.color.rgb = C_WHITE
        r.font.bold = (j == 0); r.font.name = "Arial"

    tb(slide, 400_000, H - 1_100_000, W - 500_000, 300_000,
       "次回レビュー: 2024年7月（Q2報告）", 11,
       color=RGBColor(0xCC, 0xD6, 0xE8))
    footer_band(slide, "Confidential  |  Sales Report Q1 2024  |  株式会社サンプル商事")


# ────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    # ブランクレイアウトが無い場合は追加
    while len(prs.slide_layouts) < 7:
        prs.slide_layouts._sldLayoutLst.append(
            prs.slide_layouts[0]._element.__class__()
        )

    slide_title(prs)
    slide_table(prs)
    slide_analysis(prs)
    slide_swot(prs)
    slide_summary(prs)

    prs.save(OUT)
    print(f"生成完了: {OUT}")
