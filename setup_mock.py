"""
setup_mock.py
複雑なダミー sales_data.csv と スタイリッシュな template.pptx を
data/ ディレクトリに生成するスクリプト。
"""

import os
import random
from datetime import date, timedelta

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── カラーパレット ────────────────────────────────────────
C_NAVY    = RGBColor(0x1B, 0x2A, 0x4A)   # ダークネイビー
C_TEAL    = RGBColor(0x00, 0xB4, 0xD8)   # ティール（アクセント1）
C_ORANGE  = RGBColor(0xFF, 0x6B, 0x35)   # オレンジ（アクセント2）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xF0, 0xF4, 0xF8)   # 背景薄グレー
C_DGRAY   = RGBColor(0x4A, 0x5A, 0x6A)   # 本文グレー
C_SILVER  = RGBColor(0xCC, 0xD6, 0xE0)   # 区切り線

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "data")
os.makedirs(OUTPUT_DIR, exist_ok=True)

CSV_PATH  = os.path.join(OUTPUT_DIR, "sales_data.csv")
PPTX_PATH = os.path.join(OUTPUT_DIR, "template.pptx")


# ════════════════════════════════════════════════════════
#  1. 複雑な sales_data.csv
# ════════════════════════════════════════════════════════
def create_csv():
    random.seed(2025)

    # (単価, 原価率) のタプル
    products = {
        "クラウドプラン Pro":    (120_000, 0.38),
        "クラウドプラン Basic":   (45_000, 0.42),
        "エンタープライズ SaaS":  (380_000, 0.30),
        "保守サポート契約":        (80_000, 0.50),
        "コンサルティング":       (200_000, 0.35),
        "データ分析ツール":        (95_000, 0.40),
        "セキュリティパッケージ": (150_000, 0.33),
        "ライセンス（年間）":     (250_000, 0.28),
    }
    reps = ["田中 健", "佐藤 美咲", "鈴木 一郎", "山田 花子",
            "中村 大輔", "小林 奈々", "加藤 誠", "吉田 彩",
            "渡辺 拓也", "松本 理恵"]
    regions = ["東京", "大阪", "名古屋", "福岡", "札幌", "仙台", "広島", "海外"]

    # 担当者ごとに得意地域を設定（リアリティ向上）
    rep_region_bias = {
        "田中 健":    ["東京", "東京", "大阪"],
        "佐藤 美咲":  ["東京", "仙台", "東京"],
        "鈴木 一郎":  ["大阪", "名古屋", "大阪"],
        "山田 花子":  ["大阪", "広島", "名古屋"],
        "中村 大輔":  ["福岡", "広島", "名古屋"],
        "小林 奈々":  ["名古屋", "東京", "大阪"],
        "加藤 誠":    ["札幌", "仙台", "東京"],
        "吉田 彩":    ["海外", "海外", "東京"],
        "渡辺 拓也":  ["東京", "大阪", "福岡"],
        "松本 理恵":  ["海外", "東京", "大阪"],
    }

    start = date(2025, 1, 1)
    rows = []

    for i in range(360):
        # 月ごとに売上トレンドを設定（Q1低め→Q2回復→Q3最高→Q4安定）
        offset_days = random.randint(0, 179)
        d = start + timedelta(days=offset_days)
        month = d.month
        trend = {1: 0.75, 2: 0.80, 3: 0.95, 4: 1.05, 5: 1.15, 6: 1.10}.get(month, 1.0)

        prod_name = random.choice(list(products.keys()))
        unit_price, cost_rate = products[prod_name]
        rep = random.choice(reps)
        region = random.choice(rep_region_bias[rep])

        qty = random.choices([1, 2, 3, 5, 10], weights=[40, 25, 15, 12, 8])[0]

        # 海外は単価1.3倍
        multiplier = 1.3 if region == "海外" else 1.0
        # トレンド×乱数で自然な揺れを加える
        amount = int(qty * unit_price * multiplier * trend * random.uniform(0.9, 1.1))
        # 1000円単位に丸め
        amount = round(amount / 1000) * 1000

        # 原価・利益額・利益率を計算（原価率に±5%の揺れを加える）
        actual_cost_rate = cost_rate * random.uniform(0.95, 1.05)
        cost   = round(amount * actual_cost_rate / 1000) * 1000
        profit = amount - cost
        margin = round(profit / amount * 100, 1) if amount > 0 else 0.0

        rows.append({
            "日付":       d.strftime("%Y-%m-%d"),
            "商品名":     prod_name,
            "担当者":     rep,
            "地域":       region,
            "数量":       qty,
            "売上金額":   amount,
            "原価":       cost,
            "利益額":     profit,
            "利益率(%)":  margin,
        })

    df = pd.DataFrame(rows).sort_values("日付").reset_index(drop=True)
    df.to_csv(CSV_PATH, index=False, encoding="utf-8-sig")
    xlsx_path = CSV_PATH.replace(".csv", ".xlsx")
    df.to_excel(xlsx_path, index=False)
    total  = df["売上金額"].sum()
    profit = df["利益額"].sum()
    margin = round(profit / total * 100, 1)
    print(f"[OK] CSV  生成: {CSV_PATH}  ({len(df)} 行 / 総売上 {total:,}円 / 利益率 {margin}%)")
    print(f"[OK] XLSX 生成: {xlsx_path}")


# ════════════════════════════════════════════════════════
#  2. スタイリッシュな template.pptx
# ════════════════════════════════════════════════════════

def _no_border(shape):
    """図形の枠線を透明にする。"""
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, fill_color, border=False):
    """塗りつぶし矩形を追加。"""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        _no_border(shape)
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def _set_slide_bg(slide, color: RGBColor):
    """スライド背景色を設定。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    W = 13.33  # スライド幅
    H = 7.5    # スライド高さ

    # ════════════════════════════════════════
    #  スライド 1：表紙
    # ════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    _set_slide_bg(s1, C_NAVY)

    # 左側ティールの縦帯
    _add_rect(s1, 0, 0, 0.45, H, C_TEAL)

    # 右下のオレンジ三角形っぽい大きな斜め帯（矩形を重ねて表現）
    _add_rect(s1, 8.5, 5.2, 4.83, 2.3, C_ORANGE)
    _add_rect(s1, 9.5, 5.8, 3.83, 1.7, C_NAVY)   # 上から重ねてネイビーで一部隠す

    # 右上の細いティール帯
    _add_rect(s1, 0.45, 0, W - 0.45, 0.08, C_TEAL)

    # 中央：白いセパレーター横線
    _add_rect(s1, 0.8, 3.65, 7.5, 0.05, C_TEAL)

    # タイトル
    _add_textbox(s1, "{{report_title}}",
                 left=0.8, top=1.9, width=11.5, height=1.4,
                 font_size=38, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)

    # サブタイトル
    _add_textbox(s1, "Sales Performance Report",
                 left=0.8, top=3.3, width=9, height=0.55,
                 font_size=15, bold=False, color=C_TEAL,
                 align=PP_ALIGN.LEFT, italic=True)

    # 作成日
    _add_textbox(s1, "{{report_date}}",
                 left=0.8, top=3.85, width=9, height=0.5,
                 font_size=14, color=C_SILVER, align=PP_ALIGN.LEFT)

    # 右下の会社名プレースホルダー
    _add_textbox(s1, "ACME Corporation",
                 left=9.0, top=6.7, width=4.0, height=0.5,
                 font_size=12, color=C_WHITE, align=PP_ALIGN.RIGHT, italic=True)

    # ページ番号エリア（装飾）
    _add_rect(s1, 0.45, H - 0.5, 0.35, 0.5, C_ORANGE)

    # ════════════════════════════════════════
    #  スライド 2：売上サマリー
    # ════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    _set_slide_bg(s2, C_LGRAY)

    # トップヘッダーバー（ネイビー）
    _add_rect(s2, 0, 0, W, 1.25, C_NAVY)

    # ヘッダー内 左アクセント
    _add_rect(s2, 0, 0, 0.45, 1.25, C_TEAL)

    # ヘッダータイトル
    _add_textbox(s2, "売上サマリー",
                 left=0.65, top=0.22, width=8, height=0.75,
                 font_size=28, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)

    # ヘッダー右：英語サブタイトル
    _add_textbox(s2, "SALES SUMMARY",
                 left=8.5, top=0.4, width=4.5, height=0.5,
                 font_size=13, bold=False, color=C_TEAL,
                 align=PP_ALIGN.RIGHT, italic=True)

    # 白いコンテンツカード
    _add_rect(s2, 0.45, 1.4, W - 0.65, 5.6, C_WHITE)

    # カード左の縦アクセントバー
    _add_rect(s2, 0.45, 1.4, 0.07, 5.6, C_TEAL)

    # サマリーテキスト
    _add_textbox(s2, "{{summary_text}}",
                 left=0.75, top=1.55, width=12.2, height=5.25,
                 font_size=15, color=C_DGRAY,
                 align=PP_ALIGN.LEFT)

    # フッターバー
    _add_rect(s2, 0, H - 0.35, W, 0.35, C_NAVY)
    _add_textbox(s2, "CONFIDENTIAL  |  ACME Corporation",
                 left=0.3, top=H - 0.32, width=8, height=0.3,
                 font_size=9, color=C_SILVER, align=PP_ALIGN.LEFT)
    _add_textbox(s2, "2 / 3",
                 left=12.5, top=H - 0.32, width=0.6, height=0.3,
                 font_size=9, color=C_SILVER, align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════
    #  スライド 3：所見・次月の方針
    # ════════════════════════════════════════
    s3 = prs.slides.add_slide(blank)
    _set_slide_bg(s3, C_LGRAY)

    # トップヘッダーバー（ネイビー）
    _add_rect(s3, 0, 0, W, 1.25, C_NAVY)

    # ヘッダー内 左アクセント（オレンジ）
    _add_rect(s3, 0, 0, 0.45, 1.25, C_ORANGE)

    # ヘッダータイトル
    _add_textbox(s3, "所見・次月の方針",
                 left=0.65, top=0.22, width=8, height=0.75,
                 font_size=28, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)

    # ヘッダー右：英語サブタイトル
    _add_textbox(s3, "INSIGHTS & NEXT ACTIONS",
                 left=7.5, top=0.4, width=5.5, height=0.5,
                 font_size=13, bold=False, color=C_ORANGE,
                 align=PP_ALIGN.RIGHT, italic=True)

    # 白いコンテンツカード
    _add_rect(s3, 0.45, 1.4, W - 0.65, 5.6, C_WHITE)

    # カード左の縦アクセントバー（オレンジ）
    _add_rect(s3, 0.45, 1.4, 0.07, 5.6, C_ORANGE)

    # 分析テキスト
    _add_textbox(s3, "{{analysis_text}}",
                 left=0.75, top=1.55, width=12.2, height=5.25,
                 font_size=15, color=C_DGRAY,
                 align=PP_ALIGN.LEFT)

    # フッターバー
    _add_rect(s3, 0, H - 0.35, W, 0.35, C_NAVY)
    _add_textbox(s3, "CONFIDENTIAL  |  ACME Corporation",
                 left=0.3, top=H - 0.32, width=8, height=0.3,
                 font_size=9, color=C_SILVER, align=PP_ALIGN.LEFT)
    _add_textbox(s3, "3 / 3",
                 left=12.5, top=H - 0.32, width=0.6, height=0.3,
                 font_size=9, color=C_SILVER, align=PP_ALIGN.RIGHT)

    prs.save(PPTX_PATH)
    print(f"[OK] PPTX 生成: {PPTX_PATH}  (3 スライド / モダンデザイン)")


# ── main ─────────────────────────────────────────────────
if __name__ == "__main__":
    print("=== setup_mock.py: サンプルデータ生成 ===")
    create_csv()
    create_pptx()
    print("=== 完了 ===")
