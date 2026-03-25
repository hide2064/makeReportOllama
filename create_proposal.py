"""
create_proposal.py
本システム導入提案資料 (PowerPoint) を生成するスクリプト。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "makeReportOllama_proposal.pptx")

# ── カラーパレット ──────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x37, 0x64)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xE0, 0xE7, 0xFF)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xD9, 0x77, 0x06)


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, text, left, top, width, height,
         font_size=18, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, italic=False, bg_color=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text      = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ── スライド生成関数 ──────────────────────────────────────

def slide_title(prs):
    """スライド 1: 表紙"""
    s = add_slide(prs)
    _set_bg(s, NAVY)
    _rect(s, 0, 5.8, 13.33, 1.7, INDIGO)
    _box(s, "makeReportOllama", 1, 1.2, 11.33, 1,
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _box(s, "ローカル LLM による売上報告書 自動生成システム", 1, 2.5, 11.33, 0.7,
         font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)
    _box(s, "導入提案資料", 1, 3.4, 11.33, 0.5,
         font_size=16, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
    _box(s, "2026年3月", 1, 6.5, 11.33, 0.5,
         font_size=14, color=WHITE, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    """スライド 2: アジェンダ"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "アジェンダ", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)
    items = [
        "01  課題の背景",
        "02  システム概要",
        "03  主な機能・特徴",
        "04  システム構成",
        "05  導入効果・メリット",
        "06  導入ステップ",
    ]
    for i, item in enumerate(items):
        _box(s, item, 2.5, 1.5 + i * 0.82, 8, 0.65,
             font_size=18, color=NAVY, align=PP_ALIGN.LEFT)


def slide_problem(prs):
    """スライド 3: 課題の背景"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "課題の背景", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)
    problems = [
        ("⏱  時間コスト",   "毎月の売上報告書作成に担当者が 3〜4 時間を費やしている"),
        ("✏️  手作業ミス",   "Excel からのコピー&ペーストによる転記ミスが発生しやすい"),
        ("📊  品質のばらつき", "担当者によって文章の深さ・品質に差が生じる"),
        ("🔒  情報漏洩リスク", "外部クラウド AI への売上データ送信はセキュリティリスクとなる"),
    ]
    for i, (title, desc) in enumerate(problems):
        top = 1.5 + i * 1.3
        _rect(s, 0.5, top, 12.33, 1.1, LIGHT)
        _box(s, title, 0.7, top + 0.1, 3.5, 0.4,
             font_size=15, bold=True, color=NAVY)
        _box(s, desc,  4.3, top + 0.05, 8.3, 0.9,
             font_size=14, color=GRAY, align=PP_ALIGN.LEFT)


def slide_overview(prs):
    """スライド 4: システム概要"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "システム概要", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)
    _box(s, "makeReportOllama は、Excel 売上データと PPTX テンプレートをアップロードするだけで、\n"
            "ローカル LLM (Ollama) が自動で分析・文章生成し、PowerPoint 報告書を出力するシステムです。",
         0.5, 1.4, 12.33, 1.0,
         font_size=15, color=GRAY)
    # フロー図（矢印を文字で代用）
    steps = [
        ("① Excel\nアップロード", INDIGO),
        ("② データ\n集計・分析", NAVY),
        ("③ LLM で\n文章生成", INDIGO),
        ("④ PPTX\n自動生成", NAVY),
        ("⑤ ダウンロード", GREEN),
    ]
    for i, (label, color) in enumerate(steps):
        left = 0.5 + i * 2.55
        _rect(s, left, 2.8, 2.2, 1.3, color)
        _box(s, label, left, 2.85, 2.2, 1.2,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            _box(s, "→", left + 2.2, 3.1, 0.35, 0.6,
                 font_size=20, bold=True, color=NAVY)
    _box(s, "※ すべての処理はローカル環境で完結。外部への通信は一切発生しません。",
         0.5, 4.5, 12.33, 0.6, font_size=13, color=AMBER, italic=True)


def slide_features(prs):
    """スライド 5: 主な機能・特徴"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "主な機能・特徴", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)
    features = [
        ("🔒 完全ローカル処理",  "Ollama を利用し、データが社外に出ない安全な環境で動作します。"),
        ("📂 Excel → PPTX 自動変換", "売上データをアップロードするだけで、分析・報告書生成が完結します。"),
        ("💬 AI による文章自動生成", "売上サマリーと次月方針の2種類の文章をローカル LLM が自動作成。"),
        ("🖥️ ブラウザで操作完結", "専用アプリ不要。ブラウザから直感的に操作できます。"),
        ("⚡ ローディング表示",   "CPU 推論中はリアルタイムのローディング表示で進捗を把握できます。"),
    ]
    for i, (title, desc) in enumerate(features):
        top = 1.5 + i * 1.1
        _box(s, title, 0.5, top, 4.5, 0.5,
             font_size=15, bold=True, color=NAVY)
        _box(s, desc,  5.0, top + 0.02, 8.0, 0.9,
             font_size=14, color=GRAY)
        _rect(s, 0.5, top + 0.55, 12.33, 0.02, LIGHT)


def slide_architecture(prs):
    """スライド 6: システム構成"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "システム構成", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)

    # ブラウザ層
    _rect(s, 0.5, 1.4, 5.5, 1.5, LIGHT)
    _box(s, "フロントエンド\n(React + Vite)", 0.6, 1.5, 5.3, 1.2,
         font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _box(s, "ポート: 5173", 0.6, 2.7, 5.3, 0.4,
         font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    _box(s, "HTTP POST\n/api/generate", 6.2, 1.9, 1.5, 0.8,
         font_size=11, color=INDIGO, align=PP_ALIGN.CENTER)

    # API 層
    _rect(s, 7.8, 1.4, 5.0, 1.5, RGBColor(0xC7, 0xD2, 0xFE))
    _box(s, "バックエンド\n(FastAPI + Uvicorn)", 7.9, 1.5, 4.8, 1.2,
         font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _box(s, "ポート: 8000", 7.9, 2.7, 4.8, 0.4,
         font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # サービス層
    services = ["excel_reader\n(pandas)", "ollama_client\n(httpx)", "pptx_generator\n(python-pptx)"]
    for i, svc in enumerate(services):
        _rect(s, 0.5 + i * 4.3, 3.3, 4.0, 1.3, RGBColor(0xEE, 0xF2, 0xFF))
        _box(s, svc, 0.6 + i * 4.3, 3.4, 3.8, 1.1,
             font_size=13, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # Ollama
    _rect(s, 4.3, 5.0, 4.7, 1.0, RGBColor(0xD1, 0xFA, 0xE5))
    _box(s, "Ollama  (http://localhost:11434)\nローカル LLM（CPU 推論）",
         4.4, 5.05, 4.5, 0.9,
         font_size=13, bold=True, color=RGBColor(0x06, 0x5F, 0x46), align=PP_ALIGN.CENTER)


def slide_benefits(prs):
    """スライド 7: 導入効果・メリット"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "導入効果・メリット", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)

    # Before / After
    _box(s, "導入前", 0.5, 1.4, 5.8, 0.5,
         font_size=16, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    _box(s, "導入後", 7.2, 1.4, 5.8, 0.5,
         font_size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _box(s, "→", 6.0, 1.4, 1.0, 0.5,
         font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    befores = ["報告書作成に 3〜4 時間", "手作業による転記ミスあり", "担当者ごとに品質差あり", "外部 AI への送信リスク"]
    afters  = ["自動生成で 5〜10 分以内", "自動集計でミスゼロ", "LLM による均一品質", "完全ローカル処理で安全"]
    for i, (b, a) in enumerate(zip(befores, afters)):
        top = 2.1 + i * 1.1
        _rect(s, 0.5, top, 5.8, 0.85, RGBColor(0xFF, 0xF7, 0xED))
        _box(s, f"✗  {b}", 0.6, top + 0.1, 5.6, 0.6,
             font_size=14, color=AMBER)
        _rect(s, 7.2, top, 5.8, 0.85, RGBColor(0xF0, 0xFD, 0xF4))
        _box(s, f"✓  {a}", 7.3, top + 0.1, 5.6, 0.6,
             font_size=14, color=GREEN)


def slide_steps(prs):
    """スライド 8: 導入ステップ"""
    s = add_slide(prs)
    _set_bg(s, WHITE)
    _rect(s, 0, 0, 13.33, 1.2, NAVY)
    _box(s, "導入ステップ", 0.5, 0.2, 12, 0.8,
         font_size=28, bold=True, color=WHITE)
    steps = [
        ("STEP 1", "Ollama インストール & モデルダウンロード",
         "ollama pull llama3.2  でローカル LLM を準備します。"),
        ("STEP 2", "リポジトリをクローン",
         "git clone でソースコードを取得します。"),
        ("STEP 3", "start.bat を実行",
         "ダブルクリックするだけで環境構築・サーバー起動・ブラウザ起動が完了します。"),
        ("STEP 4", "ブラウザからファイルをアップロード",
         "Excel と PPTX テンプレートを選択して「レポートを生成する」ボタンを押すだけです。"),
    ]
    for i, (step, title, desc) in enumerate(steps):
        top = 1.5 + i * 1.35
        _rect(s, 0.5, top, 1.5, 1.1, INDIGO)
        _box(s, step, 0.5, top + 0.25, 1.5, 0.6,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _box(s, title, 2.3, top + 0.05, 10.5, 0.45,
             font_size=15, bold=True, color=NAVY)
        _box(s, desc,  2.3, top + 0.5,  10.5, 0.55,
             font_size=13, color=GRAY)


def slide_close(prs):
    """スライド 9: 締め"""
    s = add_slide(prs)
    _set_bg(s, NAVY)
    _rect(s, 0, 3.2, 13.33, 2.0, INDIGO)
    _box(s, "ありがとうございました", 1, 1.5, 11.33, 1.0,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _box(s, "makeReportOllama — ローカル LLM による報告書自動生成", 1, 2.7, 11.33, 0.6,
         font_size=16, color=LIGHT, align=PP_ALIGN.CENTER)
    _box(s, "ご不明な点はいつでもお問い合わせください。",
         1, 3.6, 11.33, 0.6,
         font_size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ── main ─────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_overview(prs)
    slide_features(prs)
    slide_architecture(prs)
    slide_benefits(prs)
    slide_steps(prs)
    slide_close(prs)

    prs.save(OUTPUT_PATH)
    print(f"[OK] 提案資料を生成しました: {OUTPUT_PATH}  ({len(prs.slides)} スライド)")


if __name__ == "__main__":
    main()
