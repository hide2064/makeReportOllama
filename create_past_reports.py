"""
create_past_reports.py
過去 4 年分 (2020〜2023 Q1) のサンプル売上報告書 PPTX を一括生成する。
スライド構成は sample_report_2024Q1.pptx と同じ 5 スライド構成。
数値・分析は年ごとに変化させ、2024Q1 への自然な流れを表現する。
"""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

OUT_DIR = Path(__file__).parent / "data"

# ── カラーパレット ────────────────────────────────────────────
C_NAVY  = RGBColor(0x1B, 0x2E, 0x4C)
C_NAVY2 = RGBColor(0x2C, 0x4A, 0x7A)
C_GOLD  = RGBColor(0xC4, 0x97, 0x3E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY  = RGBColor(0x55, 0x65, 0x7A)
C_LGRAY = RGBColor(0xF0, 0xF2, 0xF5)
C_GREEN = RGBColor(0x16, 0x75, 0x3A)
C_RED   = RGBColor(0xCC, 0x28, 0x28)
C_BLUE  = RGBColor(0x1A, 0x56, 0xA0)
C_PURP  = RGBColor(0x8B, 0x5C, 0xF6)
C_TEAL  = RGBColor(0x0E, 0x9F, 0x6E)

W = 12_192_000
H =  6_858_000


# ── 描画ユーティリティ ────────────────────────────────────────
def rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line: s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, size, bold=False,
       color=C_TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Arial"
    return box


def header_band(slide, title, sub=""):
    rect(slide, 0, 0, W, 680_000, C_NAVY)
    rect(slide, 0, 0, 160_000, 680_000, C_GOLD)
    tb(slide, 220_000, 100_000, W - 500_000, 480_000,
       title, 18, bold=True, color=C_WHITE)
    if sub:
        tb(slide, W - 2_200_000, 100_000, 2_100_000, 480_000,
           sub, 9, color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.RIGHT)


def footer_band(slide, note="社外秘 — 取扱注意"):
    rect(slide, 0, H - 360_000, W, 360_000, C_NAVY)
    rect(slide, 0, H - 360_000, W, 26_000, C_GOLD)
    tb(slide, 360_000, H - 310_000, W - 720_000, 280_000,
       note, 8, color=RGBColor(0xAA, 0xB8, 0xCC))


# ── 年次データ定義 ─────────────────────────────────────────────
REPORTS = [
    {
        "filename": "sample_report_2020Q1.pptx",
        "period":   "2020年1月〜3月",
        "quarter":  "2020年 第1四半期",
        "created":  "2020年4月7日",
        "next_review": "2020年7月（Q2報告）",
        # サマリーカード
        "total":    "2,640万円",
        "total_vs": "▲5.2%",
        "total_cc": C_RED,
        "qty":      "897 個",
        "qty_vs":   "▲4.8%",
        "qty_cc":   C_RED,
        "margin":   "49.2%",
        "margin_vs":"▲1.5pt",
        "margin_cc":C_RED,
        "top_prod": "プレミアムP",
        "top_amt":  "¥1,050万",
        # 月別テーブル
        "monthly": [
            ["1月", "380万", "240万", "230万", "90万",  "940万",  "50.1%"],
            ["2月", "340万", "230万", "185万", "87万",  "842万",  "48.6%"],
            ["3月", "330万", "240万", "205万", "83万",  "858万",  "48.9%"],
            ["合計","1,050万","710万","620万", "260万", "2,640万","49.2%"],
        ],
        # 商品別
        "products": [
            ("プレミアムプラン",   "1,050万", 39.8, C_BLUE),
            ("スタンダードプラン",   "710万", 26.9, C_TEAL),
            ("コンサルティング",     "620万", 23.5, C_PURP),
            ("保守サポート",         "260万",  9.8, C_GOLD),
        ],
        # 地域別
        "regions": [
            ("東京",   "1,010万", "▲3.2%", C_RED),
            ("大阪",     "750万", "+1.1%", C_GREEN),
            ("名古屋",   "530万", "▲2.5%", C_RED),
            ("福岡",     "220万", "▲8.4%", C_RED),
            ("札幌",     "130万","▲12.1%", C_RED),
        ],
        # SWOT
        "strength_title": "強み・安定要因",
        "weakness_title":  "課題・リスク",
        "policy_title":    "今後の方針",
        "strengths": [
            "保守サポートが景気変動に左右されない安定収益を維持",
            "大阪エリアのみ前Q比プラスを維持し基盤の底堅さを証明",
            "プレミアムPの顧客離脱率が低く LTV（顧客生涯価値）は高水準",
            "既存顧客からの継続受注が全売上の 62% を占め収益を下支え",
        ],
        "weaknesses": [
            "新型コロナウイルスの感染拡大により対面商談が事実上停止",
            "コンサルティング売上が前Q比 ▲31% と大幅に落ち込む",
            "北海道・九州エリアで担当者の外出自粛による失注が多発",
            "新規顧客獲得件数が前Q比 ▲47% と急減しパイプラインが枯渇",
        ],
        "policies": [
            "オンライン商談ツール（Zoom/Teams）を全拠点に即時導入",
            "コンサルティングをオンライン完結型へ改変し早期売上回復を図る",
            "既存顧客への保守S拡張提案を強化し解約防止とアップセルを推進",
            "Q2 に向けウェビナー形式の新規リード獲得施策を月 2 回実施",
        ],
        # まとめ
        "summary_points": [
            "Q1 総売上 2,640万円 (前Q比 ▲5.2%) — コロナ禍により減収",
            "コンサルティングが対面停止で ▲31% と最大の下落要因",
            "保守サポート・大阪エリアが下支えとなり影響を最小化",
            "オンライン商談への全面移行で Q2 の早期回復を目指す",
            "既存顧客深耕とウェビナー施策で新規パイプラインを再構築する",
        ],
    },
    {
        "filename": "sample_report_2021Q1.pptx",
        "period":   "2021年1月〜3月",
        "quarter":  "2021年 第1四半期",
        "created":  "2021年4月6日",
        "next_review": "2021年7月（Q2報告）",
        "total":    "2,980万円",
        "total_vs": "+12.9%",
        "total_cc": C_GREEN,
        "qty":      "1,012 個",
        "qty_vs":   "+12.8%",
        "qty_cc":   C_GREEN,
        "margin":   "50.8%",
        "margin_vs":"+1.6pt",
        "margin_cc":C_GREEN,
        "top_prod": "プレミアムP",
        "top_amt":  "¥1,180万",
        "monthly": [
            ["1月", "400万", "255万", "270万", "75万",  "1,000万", "51.2%"],
            ["2月", "380万", "245万", "280万", "77万",   "982万",  "50.5%"],
            ["3月", "400万", "250万", "270万", "78万",   "998万",  "50.7%"],
            ["合計","1,180万","750万","820万",  "230万", "2,980万","50.8%"],
        ],
        "products": [
            ("プレミアムプラン",   "1,180万", 39.6, C_BLUE),
            ("コンサルティング",     "820万", 27.5, C_PURP),
            ("スタンダードプラン",   "750万", 25.2, C_TEAL),
            ("保守サポート",         "230万",  7.7, C_GOLD),
        ],
        "regions": [
            ("東京",   "1,140万", "+12.9%", C_GREEN),
            ("大阪",     "840万", "+12.0%", C_GREEN),
            ("名古屋",   "590万", "+11.3%", C_GREEN),
            ("福岡",     "250万", "+13.6%", C_GREEN),
            ("札幌",     "160万", "+23.1%", C_GREEN),
        ],
        "strength_title": "強み・成長要因",
        "weakness_title":  "課題・残留リスク",
        "policy_title":    "Q2 に向けた方針",
        "strengths": [
            "オンライン商談への移行が完了し全エリアでプラス成長を達成",
            "コンサルティングが前年同期比 +32.3% と力強く回復",
            "保守サポートの自動更新率が 91% に達し安定収益の柱に成長",
            "ウェビナー施策で四半期に 128 件の新規リードを獲得",
        ],
        "weaknesses": [
            "地方拠点（福岡・札幌）の担当者不足により商談キャパが不足",
            "オンライン完結型コンサルの単価が対面比 ▲15% と課題が残る",
            "スタンダードPの新機能投資が先行し利益率が 43% に低下",
            "リード獲得数は回復したが有効商談率が 38% と前年の 45% を下回る",
        ],
        "policies": [
            "Q2 に福岡・札幌各 1 名を増員し地方商談キャパを強化",
            "対面・ハイブリッド商談の再開でコンサル単価の回復を目指す",
            "スタンダードPの機能強化投資を Q2 で完了し利益率を 48% 以上に回復",
            "SDR（インサイドセールス）チーム設置でリード有効化率 45% 以上を目標",
        ],
        "summary_points": [
            "Q1 総売上 2,980万円 (前年同期比 +12.9%) — コロナ禍からV字回復",
            "全エリア・全商品でプラス成長、オンライン化が奏功",
            "コンサルティングが +32.3% で回復の牽引役に",
            "地方拠点の人員不足が次四半期の成長制約リスク",
            "Q2 増員・SDR 設置・ハイブリッド商談再開で更なる拡大を図る",
        ],
    },
    {
        "filename": "sample_report_2022Q1.pptx",
        "period":   "2022年1月〜3月",
        "quarter":  "2022年 第1四半期",
        "created":  "2022年4月5日",
        "next_review": "2022年7月（Q2報告）",
        "total":    "3,320万円",
        "total_vs": "+11.4%",
        "total_cc": C_GREEN,
        "qty":      "1,128 個",
        "qty_vs":   "+11.5%",
        "qty_cc":   C_GREEN,
        "margin":   "51.5%",
        "margin_vs":"+0.7pt",
        "margin_cc":C_GREEN,
        "top_prod": "プレミアムP",
        "top_amt":  "¥1,280万",
        "monthly": [
            ["1月", "430万", "275万", "320万", "83万",  "1,108万", "52.0%"],
            ["2月", "415万", "265万", "340万", "82万",  "1,102万", "51.2%"],
            ["3月", "435万", "270万", "320万", "85万",  "1,110万", "51.2%"],
            ["合計","1,280万","810万","980万",  "250万", "3,320万","51.5%"],
        ],
        "products": [
            ("プレミアムプラン",   "1,280万", 38.6, C_BLUE),
            ("コンサルティング",     "980万", 29.5, C_PURP),
            ("スタンダードプラン",   "810万", 24.4, C_TEAL),
            ("保守サポート",         "250万",  7.5, C_GOLD),
        ],
        "regions": [
            ("東京",   "1,265万", "+10.9%", C_GREEN),
            ("大阪",     "920万",  "+9.5%", C_GREEN),
            ("名古屋",   "655万", "+11.0%", C_GREEN),
            ("福岡",     "295万", "+18.0%", C_GREEN),
            ("札幌",     "185万", "+15.6%", C_GREEN),
        ],
        "strength_title": "競争優位・強み",
        "weakness_title":  "課題・競合リスク",
        "policy_title":    "差別化・強化方針",
        "strengths": [
            "コンサルティング売上が初めて 1,000万円に迫り収益多角化が進展",
            "福岡・札幌で増員効果が顕在化し地方商談が前Q比 +16% 超を達成",
            "顧客単価が前年比 +6.2% に向上しアップセル戦略が結果に直結",
            "解約率が 2.1% と過去最低水準を達成しカスタマーサクセスが機能",
        ],
        "weaknesses": [
            "新興クラウド競合 3 社が低価格で市場参入し見積競合案件が増加",
            "スタンダードP の値引き要求が前Q比 +28% 増加し利益率を圧迫",
            "プレミアムP の新機能開発遅延（2 ヶ月）で顧客クレームが 12 件発生",
            "採用コストが前年比 +40% に増加しオペレーションコストが上昇",
        ],
        "policies": [
            "プレミアムPの独自機能強化を Q2 内に完了し競合との差別化を明確化",
            "スタンダードPの値引き上限を設定し価格規律を徹底（ルール策定）",
            "競合対策として ROI 提案資料を整備し価格勝負から価値訴求へ転換",
            "採用コスト最適化のため内部紹介制度の報奨金を倍増し採用効率を向上",
        ],
        "summary_points": [
            "Q1 総売上 3,320万円 (前年同期比 +11.4%) — 3期連続二桁成長",
            "全エリア・全商品でプラス成長、コンサルが 980万円まで拡大",
            "競合参入で値引き圧力増大 → Q2 で差別化戦略を本格展開",
            "解約率 2.1% と過去最低を記録し顧客基盤が着実に強化",
            "Q2 は機能強化完了と価格規律徹底で収益性の一段向上を目指す",
        ],
    },
    {
        "filename": "sample_report_2023Q1.pptx",
        "period":   "2023年1月〜3月",
        "quarter":  "2023年 第1四半期",
        "created":  "2023年4月4日",
        "next_review": "2023年7月（Q2報告）",
        "total":    "3,580万円",
        "total_vs": "+7.8%",
        "total_cc": C_GREEN,
        "qty":      "1,195 個",
        "qty_vs":   "+5.9%",
        "qty_cc":   C_GREEN,
        "margin":   "52.1%",
        "margin_vs":"+0.6pt",
        "margin_cc":C_GREEN,
        "top_prod": "プレミアムP",
        "top_amt":  "¥1,420万",
        "monthly": [
            ["1月", "470万", "285万", "345万", "81万",  "1,181万", "52.4%"],
            ["2月", "475万", "290万", "355万", "79万",  "1,199万", "51.9%"],
            ["3月", "475万", "285万", "360万", "80万",  "1,200万", "51.9%"],
            ["合計","1,420万","860万","1,060万","240万", "3,580万","52.1%"],
        ],
        "products": [
            ("プレミアムプラン",   "1,420万", 39.7, C_BLUE),
            ("コンサルティング",  "1,060万", 29.6, C_PURP),
            ("スタンダードプラン",  "860万",  24.0, C_TEAL),
            ("保守サポート",        "240万",   6.7, C_GOLD),
        ],
        "regions": [
            ("東京",   "1,360万",  "+7.5%", C_GREEN),
            ("大阪",     "895万",  "+9.2%", C_GREEN),
            ("名古屋",   "680万",  "+3.8%", C_GREEN),
            ("福岡",     "430万", "+45.8%", C_GREEN),
            ("札幌",     "215万", "+16.2%", C_GREEN),
        ],
        "strength_title": "成長ドライバー・強み",
        "weakness_title":  "課題・コスト圧力",
        "policy_title":    "持続成長に向けた方針",
        "strengths": [
            "コンサルティングが初の四半期 1,000万円超えを達成し二柱体制が確立",
            "福岡エリアが +45.8% と突出した成長を遂げ地方展開戦略の有効性を実証",
            "全エリアでプラス成長を達成し地域リスク分散が完成段階に",
            "プレミアムP の更新率が 94% と過去最高を記録し優良顧客層を確保",
        ],
        "weaknesses": [
            "人件費・採用コストが前年比 +22% 増加し営業利益を圧迫",
            "保守サポート単価が競合の低価格攻勢で Q1 に ▲8% 低下",
            "急成長中の福岡担当 2 名に業務が集中し属人化リスクが顕在化",
            "スタンダードPの競合との機能差が縮小し乗り換え検討顧客が増加",
        ],
        "policies": [
            "保守サポートの価格体系を再設計し Q2 に値上げ交渉を順次実施",
            "福岡に 1 名追加採用し属人化リスクを分散（Q2 中に着任予定）",
            "スタンダードPの AI 機能追加を Q3 にリリースし差別化を再強化",
            "全商品でロールアップ型年間契約を推進し Q2 の売上変動を平準化",
        ],
        "summary_points": [
            "Q1 総売上 3,580万円 (前年同期比 +7.8%) — 4期連続増収",
            "コンサルティング 1,060万円、四半期初の 1,000万円超えを達成",
            "福岡 +45.8% と地方展開戦略が結実、全エリアプラス成長",
            "人件費増・保守S単価下落がコスト圧力として顕在化",
            "保守S値上げ・AI 機能強化・福岡増員で Q2 以降の収益性向上を図る",
        ],
    },
]


# ── スライド生成関数 ──────────────────────────────────────────
def make_slide1_title(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, C_NAVY)
    rect(slide, 0, 0, 220_000, H, C_GOLD)
    rect(slide, 220_000, H // 2 - 10_000, W - 220_000, 20_000, C_GOLD)
    tb(slide, 400_000, 1_200_000, W - 500_000, 1_000_000,
       d["quarter"], 16, color=RGBColor(0xCC, 0xD6, 0xE8))
    tb(slide, 400_000, 2_000_000, W - 500_000, 1_200_000,
       "売上報告書", 42, bold=True, color=C_WHITE)
    tb(slide, 400_000, 3_200_000, W - 500_000, 600_000,
       d["period"] + " 実績", 16, color=RGBColor(0xCC, 0xD6, 0xE8))
    tb(slide, 400_000, H - 1_400_000, W - 500_000, 400_000,
       "株式会社サンプル商事　営業企画部", 11,
       color=RGBColor(0xAA, 0xB8, 0xCC))
    tb(slide, 400_000, H - 1_000_000, W - 500_000, 400_000,
       "作成日: " + d["created"], 10,
       color=RGBColor(0xAA, 0xB8, 0xCC))
    footer_band(slide, f"Confidential  |  Sales Report {d['quarter']}")


def make_slide2_table(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, RGBColor(0xF7, 0xF8, 0xFA))
    header_band(slide, f"売上実績（{d['quarter']}）", "単位: 万円")

    BOX_T = 780_000; BOX_H = 700_000
    boxes = [
        ("総売上",       d["total"],    d["total_vs"],  d["total_cc"]),
        ("総販売数量",   d["qty"],      d["qty_vs"],    d["qty_cc"]),
        ("平均利益率",   d["margin"],   d["margin_vs"], d["margin_cc"]),
        ("最高売上商品", d["top_prod"], d["top_amt"],   C_BLUE),
    ]
    bw = (W - 800_000) // 4
    for i, (lbl, val, sub, sc) in enumerate(boxes):
        bx = 400_000 + i * (bw + 80_000)
        rect(slide, bx, BOX_T, bw, BOX_H, C_WHITE)
        rect(slide, bx, BOX_T, bw, 12_000, sc)
        tb(slide, bx + 80_000, BOX_T + 60_000, bw - 100_000, 260_000,
           lbl, 9, color=C_GRAY)
        tb(slide, bx + 80_000, BOX_T + 280_000, bw - 100_000, 320_000,
           val, 15, bold=True, color=C_TEXT)
        tb(slide, bx + 80_000, BOX_T + 560_000, bw - 100_000, 180_000,
           sub, 9, bold=True, color=sc)

    TBL_T = 1_620_000
    headers = ["月", "プレミアムP", "スタンダードP", "コンサルティング", "保守S", "合計", "利益率"]
    col_w = [320_000, 1_380_000, 1_380_000, 1_700_000, 900_000, 1_150_000, 900_000]
    row_h = 440_000

    for ci, (hdr, cw) in enumerate(zip(headers, col_w)):
        cx = 400_000 + sum(col_w[:ci])
        rect(slide, cx, TBL_T, cw, row_h, C_NAVY2)
        tb(slide, cx + 30_000, TBL_T + 80_000, cw - 40_000, row_h - 80_000,
           hdr, 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(d["monthly"]):
        ry = TBL_T + row_h * (ri + 1)
        bg = RGBColor(0xEE, 0xF3, 0xFA) if ri % 2 == 0 else C_WHITE
        if ri == 3: bg = RGBColor(0xE8, 0xEE, 0xF4)
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            cx = 400_000 + sum(col_w[:ci])
            rect(slide, cx, ry, cw, row_h, bg)
            fc = C_GREEN if (ri == 3 and ci == 5) else C_TEXT
            if cell.startswith("▲"): fc = C_RED
            tb(slide, cx + 30_000, ry + 80_000, cw - 40_000, row_h - 80_000,
               cell, 9.5, bold=(ri == 3), color=fc, align=PP_ALIGN.CENTER)

    footer_band(slide)


def make_slide3_analysis(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, RGBColor(0xF7, 0xF8, 0xFA))
    header_band(slide, "売上分析", d["quarter"])

    COL_T = 820_000; COL_H = H - COL_T - 400_000
    LW = (W - 800_000) // 2 - 60_000
    RX = 400_000 + LW + 120_000

    # 左: 商品別
    rect(slide, 400_000, COL_T, LW, COL_H, C_WHITE)
    rect(slide, 400_000, COL_T, LW, 380_000, C_NAVY2)
    tb(slide, 520_000, COL_T + 80_000, LW, 300_000,
       "● 商品別売上構成", 11, bold=True, color=C_WHITE)

    BY = COL_T + 420_000
    for name, amt, pct, color in d["products"]:
        bar_w = max(int(LW * 0.80 * pct / 45), 50_000)
        rect(slide, 520_000, BY, bar_w, 160_000, color)
        tb(slide, 520_000, BY + 170_000, LW - 80_000, 200_000,
           f"{name}  {amt} ({pct}%)", 8.5, color=C_TEXT)
        BY += 380_000

    # 右: 地域別
    rect(slide, RX, COL_T, LW, COL_H, C_WHITE)
    rect(slide, RX, COL_T, LW, 380_000, C_GOLD)
    tb(slide, RX + 120_000, COL_T + 80_000, LW, 300_000,
       "● 地域別売上 & 前Q比", 11, bold=True, color=C_NAVY)

    RY = COL_T + 420_000
    for region, amt, chg, cc in d["regions"]:
        rect(slide, RX + 120_000, RY + 20_000, LW - 200_000, 200_000, C_LGRAY)
        tb(slide, RX + 160_000, RY + 40_000, 800_000, 180_000,
           region, 10, bold=True, color=C_TEXT)
        tb(slide, RX + 900_000, RY + 40_000, 700_000, 180_000,
           amt, 10, color=C_TEXT, align=PP_ALIGN.RIGHT)
        tb(slide, RX + 1_620_000, RY + 40_000, 600_000, 180_000,
           chg, 10, bold=True, color=cc, align=PP_ALIGN.RIGHT)
        RY += 280_000

    footer_band(slide)


def make_slide4_swot(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, RGBColor(0xF7, 0xF8, 0xFA))
    header_band(slide, "現状分析：強み・課題と今後の方針")

    CARD_T = 800_000; CARD_H = H - CARD_T - 500_000
    CW = (W - 1_000_000) // 3
    icons = ["◎", "▲", "→"]
    cards = [
        (d["strength_title"], C_GREEN, d["strengths"], icons[0]),
        (d["weakness_title"],  C_RED,  d["weaknesses"], icons[1]),
        (d["policy_title"],   C_BLUE,  d["policies"],   icons[2]),
    ]

    for i, (title, color, bullets, icon) in enumerate(cards):
        cx = 400_000 + i * (CW + 100_000)
        rect(slide, cx, CARD_T, CW, CARD_H, C_WHITE)
        rect(slide, cx, CARD_T, CW, 400_000, color)
        tb(slide, cx + 80_000, CARD_T + 60_000, CW - 100_000, 300_000,
           f"{icon}  {title}", 11, bold=True, color=C_WHITE)
        box = slide.shapes.add_textbox(
            Emu(cx + 80_000), Emu(CARD_T + 460_000),
            Emu(CW - 160_000), Emu(CARD_H - 500_000)
        )
        tf = box.text_frame; tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            r = p.add_run(); r.text = f"• {b}"
            r.font.size = Pt(9.5); r.font.color.rgb = C_TEXT
            r.font.name = "Arial"

    footer_band(slide)


def make_slide5_summary(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, C_NAVY)
    rect(slide, 0, 0, 220_000, H, C_GOLD)
    rect(slide, 220_000, H // 2 - 10_000, W - 220_000, 20_000, C_GOLD)
    tb(slide, 400_000, 900_000, W - 500_000, 600_000,
       "まとめ", 28, bold=True, color=C_GOLD)

    box = slide.shapes.add_textbox(
        Emu(400_000), Emu(1_600_000), Emu(W - 600_000), Emu(3_600_000)
    )
    tf = box.text_frame; tf.word_wrap = True
    for j, pt in enumerate(d["summary_points"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run(); r.text = f"  ✓  {pt}"
        r.font.size = Pt(12); r.font.color.rgb = C_WHITE
        r.font.bold = (j == 0); r.font.name = "Arial"

    tb(slide, 400_000, H - 1_100_000, W - 500_000, 300_000,
       "次回レビュー: " + d["next_review"], 11,
       color=RGBColor(0xCC, 0xD6, 0xE8))
    footer_band(slide,
       f"Confidential  |  Sales Report {d['quarter']}  |  株式会社サンプル商事")


# ── メイン ────────────────────────────────────────────────────
def build_report(d: dict):
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    make_slide1_title(prs, d)
    make_slide2_table(prs, d)
    make_slide3_analysis(prs, d)
    make_slide4_swot(prs, d)
    make_slide5_summary(prs, d)

    out = OUT_DIR / d["filename"]
    prs.save(out)
    return out


if __name__ == "__main__":
    for d in REPORTS:
        out = build_report(d)
        print(f"生成: {out.name}")
    print(f"\n{len(REPORTS)} 件の報告書を {OUT_DIR} に出力しました。")
