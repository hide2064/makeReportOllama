"""
create_annual_reports.py
sample_complex.csv の実データをベースに 2020〜2024 年の年次売上報告書 PPTX を一括生成する。
"""

import csv
import io
import sys
from collections import defaultdict
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

try:
    import matplotlib
    matplotlib.use("Agg")
    matplotlib.rcParams.update({
        "font.family":     "sans-serif",
        "font.sans-serif": ["Meiryo", "MS Gothic", "Yu Gothic", "DejaVu Sans"],
        "axes.unicode_minus": False,
    })
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mtick
    _MPL_OK = True
except ImportError:
    _MPL_OK = False

# ── パス設定 ──────────────────────────────────────────────────────
BASE_DIR = Path(__file__).parent
DATA_DIR = BASE_DIR / "data"
OUT_DIR  = DATA_DIR   # 既存サンプルと同じ場所に出力

# ── カラーパレット ────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x12, 0x28)
NAVY2  = RGBColor(0x16, 0x32, 0x60)
GOLD   = RGBColor(0xD4, 0x94, 0x1A)
TEAL   = RGBColor(0x00, 0xA3, 0x9A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF6, 0xF8, 0xFC)
TEXT   = RGBColor(0x0A, 0x12, 0x28)
GRAY   = RGBColor(0x55, 0x65, 0x7A)
GREEN  = RGBColor(0x05, 0x7A, 0x55)
RED    = RGBColor(0xCC, 0x28, 0x28)

W = 12_192_000
H =  6_858_000


# ── 描画ヘルパー ─────────────────────────────────────────────────
def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, size, bold=False,
       color=TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Arial"
    return box


def header(slide, title, sub=""):
    rect(slide, 0, 0, W, 700_000, NAVY2)
    rect(slide, 0, 0, 180_000, 700_000, GOLD)
    tb(slide, 260_000, 110_000, W - 500_000, 500_000,
       title, 18, bold=True, color=WHITE)
    if sub:
        tb(slide, W - 2_300_000, 110_000, 2_200_000, 500_000,
           sub, 9, color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.RIGHT)


def footer(slide, note="社外秘 — 取扱注意"):
    rect(slide, 0, H - 360_000, W, 360_000, NAVY)
    rect(slide, 0, H - 360_000, W, 26_000, GOLD)
    tb(slide, 360_000, H - 310_000, W - 720_000, 260_000,
       note, 8, color=RGBColor(0xAA, 0xB8, 0xCC))


def blank_slide(prs):
    for layout in prs.slide_layouts:
        if not layout.placeholders:
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── データ集計 ──────────────────────────────────────────────────
def load_data(csv_path: Path) -> dict:
    result = {}
    for y in range(2020, 2025):
        result[y] = {
            "revenue": 0, "cost": 0, "profit": 0, "count": 0,
            "monthly_rev": defaultdict(int),
            "monthly_profit": defaultdict(int),
            "cat": defaultdict(int),
            "region": defaultdict(int),
            "person": defaultdict(int),
            "ctype": defaultdict(int),
            "channel": defaultdict(int),
            "product": defaultdict(lambda: {"rev": 0, "profit": 0, "cnt": 0}),
        }

    with open(csv_path, encoding="utf-8-sig") as f:
        for row in csv.DictReader(f):
            dt = datetime.strptime(row["日付"], "%Y-%m-%d")
            y = dt.year
            if y not in result:
                continue
            d = result[y]
            rev = int(row["売上金額"])
            pro = int(row["利益額"])
            d["revenue"] += rev
            d["cost"]    += int(row["原価"])
            d["profit"]  += pro
            d["count"]   += 1
            d["monthly_rev"][dt.month]    += rev
            d["monthly_profit"][dt.month] += pro
            d["cat"][row["カテゴリ"]]     += rev
            d["region"][row["地域"]]       += rev
            d["person"][row["担当者"]]     += rev
            d["ctype"][row["顧客タイプ"]] += rev
            d["channel"][row["チャネル"]] += rev
            pname = row["商品名"]
            d["product"][pname]["rev"]    += rev
            d["product"][pname]["profit"] += pro
            d["product"][pname]["cnt"]    += int(row["数量"])

    return result


def fmt_m(yen: int) -> str:
    """円 → M円 表示"""
    return f"{yen / 1_000_000:.1f}M円"


def pct_diff(current: int, prev: int) -> str:
    if prev == 0:
        return "N/A"
    diff = (current - prev) / prev * 100
    sign = "+" if diff >= 0 else "▲"
    return f"{sign}{abs(diff):.1f}%"


# ── スライド生成関数群 ────────────────────────────────────────────

def slide_cover(prs, year: int):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, 0, 0, W, 26_000, GOLD)
    rect(slide, 0, H - 26_000, W, 26_000, GOLD)
    rect(slide, 0, 0, 220_000, H, NAVY2)
    # 会社名
    tb(slide, 350_000, 1_200_000, W - 700_000, 500_000,
       "株式会社サンプルコーポレーション", 14, color=RGBColor(0xCC, 0xD6, 0xE8))
    # タイトル
    tb(slide, 350_000, 1_900_000, W - 700_000, 900_000,
       f"{year}年度  年次売上報告書", 34, bold=True, color=WHITE)
    # サブタイトル
    tb(slide, 350_000, 2_900_000, W - 700_000, 500_000,
       f"対象期間: {year}年1月1日 〜 {year}年12月31日", 14,
       color=RGBColor(0xB0, 0xC4, 0xDE))
    # 作成日
    created = f"{year + 1}年1月15日"
    tb(slide, 350_000, H - 1_400_000, W - 700_000, 400_000,
       f"作成日: {created}　　機密区分: 社外秘", 11,
       color=RGBColor(0x88, 0x99, 0xAA))
    # アクセント線
    rect(slide, 350_000, 1_750_000, 1_800_000, 40_000, GOLD)


def slide_summary(prs, year: int, d: dict, prev_d: dict | None):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, LGRAY)
    header(slide, "エグゼクティブサマリー", f"{year}年度 年次報告")
    footer(slide)

    rev = d["revenue"]
    pro = d["profit"]
    pro_rate = pro / rev * 100 if rev else 0
    cnt = d["count"]

    vs_rev = pct_diff(rev, prev_d["revenue"]) if prev_d else "初年度"
    vs_pro = pct_diff(pro, prev_d["profit"])  if prev_d else "初年度"
    vs_cnt = pct_diff(cnt, prev_d["count"])   if prev_d else "初年度"

    cards = [
        ("総売上高", fmt_m(rev), vs_rev),
        ("総利益額", fmt_m(pro), vs_pro),
        ("利益率",   f"{pro_rate:.1f}%", ""),
        ("受注件数", f"{cnt}件",         vs_cnt),
    ]

    card_w = 2_700_000
    card_h = 1_600_000
    gap    =   130_000
    top    = 1_000_000
    start_l = (W - (card_w * 4 + gap * 3)) // 2

    for i, (label, val, vs) in enumerate(cards):
        l = start_l + i * (card_w + gap)
        rect(slide, l, top, card_w, card_h, WHITE)
        rect(slide, l, top, card_w, 80_000, GOLD)
        tb(slide, l + 120_000, top + 140_000, card_w - 200_000, 400_000,
           label, 11, color=GRAY)
        tb(slide, l + 120_000, top + 550_000, card_w - 200_000, 600_000,
           val, 22, bold=True, color=NAVY)
        if vs:
            c = GREEN if vs.startswith("+") else (RED if vs.startswith("▲") else GRAY)
            tb(slide, l + 120_000, top + 1_200_000, card_w - 200_000, 280_000,
               f"前年比 {vs}", 10, color=c)

    # 所見テキスト
    top2 = top + card_h + 320_000
    rect(slide, start_l, top2, card_w * 4 + gap * 3, 1_750_000, WHITE)
    rect(slide, start_l, top2, card_w * 4 + gap * 3, 70_000, NAVY2)

    yoy_rev = f"前年比{vs_rev}" if prev_d else "基準年度"
    insight = _build_insight(year, d, prev_d, yoy_rev)
    tb(slide, start_l + 180_000, top2 + 150_000,
       card_w * 4 + gap * 3 - 360_000, 1_450_000,
       insight, 10, color=TEXT)


def _build_insight(year: int, d: dict, prev_d: dict | None, yoy_rev: str) -> str:
    top_cat    = max(d["cat"],    key=d["cat"].get)
    top_region = max(d["region"], key=d["region"].get)
    top_person = max(d["person"], key=d["person"].get)
    top_prod   = max(d["product"], key=lambda k: d["product"][k]["rev"])

    cat_share = d["cat"][top_cat] / d["revenue"] * 100 if d["revenue"] else 0

    lines = [
        f"【{year}年度 総括】",
        f"年間総売上高は{fmt_m(d['revenue'])}（{yoy_rev}）、"
        f"利益率は{d['profit']/d['revenue']*100:.1f}%を達成。",
        f"最大カテゴリは「{top_cat}」で売上全体の{cat_share:.0f}%を占める。",
        f"地域別トップは「{top_region}」、主力商品は「{top_prod}」が牽引。",
        f"MVP担当者は「{top_person}」。",
    ]
    if year == 2020:
        lines.append("コロナ禍による市場収縮の影響を受け、前年比で売上が減少。"
                     "オンライン・Web チャネルへの移行が加速し、下期より回復基調に転じた。")
    elif year == 2021:
        lines.append("経済活動再開に伴い売上が回復。SaaS・クラウド需要が急拡大し、"
                     "エンタープライズプランの新規獲得が特に好調だった。")
    elif year == 2022:
        lines.append("ロシア・ウクライナ情勢に伴うコスト上昇圧力があったものの、"
                     "価格改定と生産性向上策により利益率を維持。SaaS ARR が初めて200M円を突破。")
    elif year == 2023:
        lines.append("AI・DX 投資の高まりを受けてコンサルティング需要が急増。"
                     "大手企業向けエンタープライズ案件が増加し、客単価が向上。年末商戦も好調。")
    elif year == 2024:
        lines.append("過去最高売上を更新。直販チャネルの強化と新規地域開拓が奏功。"
                     "プレミアムサポートの契約継続率が向上し、ストック型収益が安定。")
    return "\n".join(lines)


def slide_monthly_chart(prs, year: int, d: dict, prev_d: dict | None):
    if not _MPL_OK:
        return

    months = list(range(1, 13))
    rev_vals  = [d["monthly_rev"].get(m, 0)    / 1_000_000 for m in months]
    pro_vals  = [d["monthly_profit"].get(m, 0) / 1_000_000 for m in months]
    prev_vals = None
    if prev_d:
        prev_vals = [prev_d["monthly_rev"].get(m, 0) / 1_000_000 for m in months]

    fig, ax = plt.subplots(figsize=(11.0, 4.2), dpi=150)
    fig.patch.set_facecolor("#F6F8FC")
    ax.set_facecolor("#F6F8FC")

    x = range(len(months))
    bar_w = 0.38
    bars = ax.bar([i - bar_w / 2 for i in x], rev_vals, width=bar_w,
                  color="#1632A0", alpha=0.85, label=f"{year}年 売上")
    ax.bar([i + bar_w / 2 for i in x], pro_vals, width=bar_w,
           color="#00A39A", alpha=0.75, label=f"{year}年 利益")
    if prev_vals:
        ax.plot(x, prev_vals, color="#D4941A", linewidth=2,
                marker="o", markersize=4, label=f"{year - 1}年 売上（参考）")

    ax.set_xticks(list(x))
    ax.set_xticklabels([f"{m}月" for m in months], fontsize=8)
    ax.yaxis.set_major_formatter(mtick.FuncFormatter(lambda v, _: f"{v:.0f}M"))
    ax.set_ylabel("百万円", fontsize=8)
    ax.set_title(f"{year}年 月次売上・利益推移", fontsize=11, fontweight="bold",
                 color="#0A1228", pad=10)
    ax.legend(fontsize=8, loc="upper left")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", linestyle="--", alpha=0.4)

    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)

    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, LGRAY)
    header(slide, "月次売上・利益推移", f"{year}年度")
    footer(slide)
    slide.shapes.add_picture(buf, Emu(320_000), Emu(780_000),
                              Emu(W - 640_000), Emu(H - 1_350_000))


def slide_category_chart(prs, year: int, d: dict):
    if not _MPL_OK:
        return

    cat_items = sorted(d["cat"].items(), key=lambda x: -x[1])
    labels = [k for k, _ in cat_items]
    values = [v / 1_000_000 for _, v in cat_items]
    colors_pie = ["#1632A0", "#00A39A", "#D4941A", "#CC2828", "#8B5CF6", "#057A55"]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.0, 4.0), dpi=150)
    fig.patch.set_facecolor("#F6F8FC")

    # 円グラフ
    ax1.set_facecolor("#F6F8FC")
    wedges, texts, autotexts = ax1.pie(
        values, labels=labels, autopct="%1.1f%%",
        colors=colors_pie[:len(labels)], startangle=90,
        textprops={"fontsize": 8}
    )
    ax1.set_title(f"カテゴリ別売上構成比 ({year}年)", fontsize=9, fontweight="bold",
                  color="#0A1228")

    # 棒グラフ
    ax2.set_facecolor("#F6F8FC")
    bars = ax2.barh(labels[::-1], values[::-1], color=colors_pie[:len(labels)][::-1],
                    alpha=0.85)
    ax2.xaxis.set_major_formatter(mtick.FuncFormatter(lambda v, _: f"{v:.0f}M"))
    ax2.set_title(f"カテゴリ別売上金額 ({year}年)", fontsize=9, fontweight="bold",
                  color="#0A1228")
    ax2.spines[["top", "right"]].set_visible(False)
    for bar, val in zip(bars, values[::-1]):
        ax2.text(val + 0.5, bar.get_y() + bar.get_height() / 2,
                 f"{val:.1f}M", va="center", fontsize=8)
    ax2.grid(axis="x", linestyle="--", alpha=0.4)

    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)

    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, LGRAY)
    header(slide, "カテゴリ別売上分析", f"{year}年度")
    footer(slide)
    slide.shapes.add_picture(buf, Emu(320_000), Emu(780_000),
                              Emu(W - 640_000), Emu(H - 1_350_000))


def slide_product_table(prs, year: int, d: dict, prev_d: dict | None):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, LGRAY)
    header(slide, "商品別売上ランキング", f"{year}年度 Top 10")
    footer(slide)

    products = sorted(d["product"].items(), key=lambda x: -x[1]["rev"])[:10]

    col_x   = [320_000, 2_600_000, 5_100_000, 7_400_000, 9_600_000, 11_000_000]
    col_w   = [2_200_000, 2_400_000, 2_200_000, 2_100_000, 1_300_000, 1_050_000]
    headers = ["商品名", "売上金額", "利益額", "利益率", "前年比", "件数"]
    row_h   = 410_000
    top0    = 800_000

    # ヘッダー行
    rect(slide, col_x[0], top0, W - col_x[0] - 200_000, row_h, NAVY2)
    for i, h in enumerate(headers):
        tb(slide, col_x[i] + 80_000, top0 + 80_000,
           col_w[i] - 100_000, row_h - 100_000,
           h, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for rank, (pname, pv) in enumerate(products):
        row_top = top0 + row_h + rank * row_h
        bg = WHITE if rank % 2 == 0 else LGRAY
        rect(slide, col_x[0], row_top, W - col_x[0] - 200_000, row_h, bg)

        rev    = pv["rev"]
        pro    = pv["profit"]
        rate   = pro / rev * 100 if rev else 0
        cnt    = pv["cnt"]

        vs = ""
        if prev_d and pname in prev_d["product"]:
            prev_rev = prev_d["product"][pname]["rev"]
            vs = pct_diff(rev, prev_rev)

        vals = [
            (pname,             PP_ALIGN.LEFT),
            (fmt_m(rev),        PP_ALIGN.RIGHT),
            (fmt_m(pro),        PP_ALIGN.RIGHT),
            (f"{rate:.1f}%",    PP_ALIGN.CENTER),
            (vs,                PP_ALIGN.CENTER),
            (str(cnt),          PP_ALIGN.CENTER),
        ]
        for i, (txt, align) in enumerate(vals):
            c = TEXT
            if i == 4 and vs:
                c = GREEN if vs.startswith("+") else (RED if vs.startswith("▲") else GRAY)
            tb(slide, col_x[i] + 80_000, row_top + 80_000,
               col_w[i] - 100_000, row_h - 100_000,
               txt, 9, color=c, align=align)


def slide_region_person(prs, year: int, d: dict):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, LGRAY)
    header(slide, "地域別・担当者別売上", f"{year}年度")
    footer(slide)

    # 地域
    regions = sorted(d["region"].items(), key=lambda x: -x[1])
    total_rev = d["revenue"]

    PANEL_TOP = 820_000
    PANEL_H   = 5_300_000
    HALF_W    = W // 2 - 60_000

    rect(slide, 120_000, PANEL_TOP, HALF_W, PANEL_H, WHITE)
    rect(slide, 120_000, PANEL_TOP, HALF_W, 80_000, TEAL)
    tb(slide, 200_000, PANEL_TOP + 120_000, HALF_W - 160_000, 400_000,
       "地域別売上", 12, bold=True, color=NAVY)

    for i, (reg, rev) in enumerate(regions[:8]):
        y_pos = PANEL_TOP + 600_000 + i * 560_000
        share = rev / total_rev * 100 if total_rev else 0
        bar_max = HALF_W - 500_000
        bar_len = int(bar_max * share / 100)
        rect(slide, 250_000, y_pos + 260_000, bar_max, 160_000,
             RGBColor(0xE8, 0xEC, 0xF4))
        if bar_len > 0:
            rect(slide, 250_000, y_pos + 260_000, bar_len, 160_000, TEAL)
        tb(slide, 200_000, y_pos, HALF_W - 250_000, 270_000,
           reg, 10, color=TEXT)
        tb(slide, 200_000, y_pos, HALF_W - 250_000, 270_000,
           fmt_m(rev), 10, color=NAVY, align=PP_ALIGN.RIGHT)

    # 担当者
    persons = sorted(d["person"].items(), key=lambda x: -x[1])
    px = W // 2 + 60_000
    rect(slide, px, PANEL_TOP, HALF_W, PANEL_H, WHITE)
    rect(slide, px, PANEL_TOP, HALF_W, 80_000, GOLD)
    tb(slide, px + 80_000, PANEL_TOP + 120_000, HALF_W - 160_000, 400_000,
       "担当者別売上 Top 8", 12, bold=True, color=NAVY)

    for i, (person, rev) in enumerate(persons[:8]):
        y_pos = PANEL_TOP + 600_000 + i * 560_000
        share = rev / total_rev * 100 if total_rev else 0
        bar_max = HALF_W - 500_000
        bar_len = int(bar_max * share / 100)
        rect(slide, px + 130_000, y_pos + 260_000, bar_max, 160_000,
             RGBColor(0xE8, 0xEC, 0xF4))
        if bar_len > 0:
            rect(slide, px + 130_000, y_pos + 260_000, bar_len, 160_000, GOLD)
        tb(slide, px + 80_000, y_pos, HALF_W - 200_000, 270_000,
           person, 10, color=TEXT)
        tb(slide, px + 80_000, y_pos, HALF_W - 200_000, 270_000,
           fmt_m(rev), 10, color=NAVY, align=PP_ALIGN.RIGHT)


# ── メイン ────────────────────────────────────────────────────────
def build_report(year: int, d: dict, prev_d: dict | None, out_path: Path):
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    slide_cover(prs, year)
    slide_summary(prs, year, d, prev_d)
    slide_monthly_chart(prs, year, d, prev_d)
    slide_category_chart(prs, year, d)
    slide_product_table(prs, year, d, prev_d)
    slide_region_person(prs, year, d)

    prs.save(out_path)
    print(f"  保存: {out_path.name}  ({len(prs.slides)}スライド)")


def main():
    csv_path = DATA_DIR / "sample_complex.csv"
    if not csv_path.exists():
        print(f"ERROR: {csv_path} が見つかりません", file=sys.stderr)
        sys.exit(1)

    print("データ読み込み中...")
    data = load_data(csv_path)

    print(f"レポート生成開始 (2020〜2024年)")
    prev_d = None
    for year in range(2020, 2025):
        d = data[year]
        out_path = OUT_DIR / f"annual_report_{year}.pptx"
        print(f"\n[{year}年度]  売上={fmt_m(d['revenue'])}  利益率={d['profit']/d['revenue']*100:.1f}%  件数={d['count']}")
        build_report(year, d, prev_d, out_path)
        prev_d = d

    print("\n完了。生成ファイル:")
    for year in range(2020, 2025):
        p = OUT_DIR / f"annual_report_{year}.pptx"
        print(f"  {p}")


if __name__ == "__main__":
    main()
