"""
create_samples_advanced.py
コンサルタント風テンプレート PPTX と多年度サンプル売上データ (Excel) を生成する。

生成ファイル:
  data/sample_advanced.xlsx      — 2022〜2024年の売上・原価・利益率データ
  data/template_consultant.pptx  — プロフェッショナル向けデザインテンプレート
"""

import random
from datetime import date, timedelta
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

random.seed(42)

DATA_DIR = Path(__file__).parent / "data"
DATA_DIR.mkdir(exist_ok=True)

# ── カラーパレット ────────────────────────────────────────────
C_NAVY      = RGBColor(0x1B, 0x2E, 0x4C)   # ダークネイビー
C_NAVY_MID  = RGBColor(0x2C, 0x4A, 0x7A)   # ミディアムネイビー
C_GOLD      = RGBColor(0xC4, 0x97, 0x3E)   # ゴールド
C_GOLD_LIGHT= RGBColor(0xF0, 0xD9, 0x9A)   # ライトゴールド
C_BG        = RGBColor(0xF7, 0xF8, 0xFA)   # 背景（オフホワイト）
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT      = RGBColor(0x1A, 0x1A, 0x2E)   # 本文テキスト
C_TEXT_LIGHT= RGBColor(0x55, 0x65, 0x7A)   # サブテキスト
C_DIVIDER   = RGBColor(0xD8, 0xDE, 0xE6)   # 区切り線
C_LEFT_BG   = RGBColor(0xEE, 0xF3, 0xFA)   # 左列背景（薄ブルー）
C_RIGHT_BG  = RGBColor(0xFD, 0xF6, 0xEC)   # 右列背景（薄ゴールド）

# ── スライドサイズ (16:9, 33.87cm x 19.05cm) ─────────────────
W = 12_192_000   # EMU: 33.87cm
H =  6_858_000   # EMU: 19.05cm


def rgb_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, line=False):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    rgb_fill(shp, color)
    if not line:
        shp.line.fill.background()
    return shp


def add_textbox(slide, l, t, w, h, text, size, bold=False,
                color=C_TEXT, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb


# ────────────────────────────────────────────────────────────
#  テンプレート PPTX 生成
# ────────────────────────────────────────────────────────────
def create_template():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    blank_layout = prs.slide_layouts[6]   # 完全ブランク
    slide = prs.slides.add_slide(blank_layout)

    # ── 背景 ────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, H, C_BG)

    # ── ヘッダーバー（ネイビー, 上部） ───────────────────────
    HDR_H = 700_000   # ~1.9cm
    add_rect(slide, 0, 0, W, HDR_H, C_NAVY)

    # ヘッダー内: 左ロゴ帯（ゴールド縦線 + テキスト）
    add_rect(slide, 0, 0, 180_000, HDR_H, C_GOLD)
    add_textbox(slide, 220_000, 80_000, 3_000_000, HDR_H - 100_000,
                "CONFIDENTIAL  |  Sales Report", 9,
                color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.LEFT)

    # ヘッダー内: 右側 — 作成日
    add_textbox(slide, W - 2_400_000, 80_000, 2_300_000, HDR_H - 100_000,
                "{{report_date}}", 9,
                color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.RIGHT)

    # ── タイトルセクション ────────────────────────────────────
    TITLE_T = HDR_H + 160_000
    TITLE_H = 820_000
    # ゴールドアクセントライン（左端縦線）
    add_rect(slide, 360_000, TITLE_T, 55_000, TITLE_H, C_GOLD)
    # タイトル
    add_textbox(slide, 470_000, TITLE_T, W - 900_000, TITLE_H,
                "{{report_title}}", 22, bold=True,
                color=C_NAVY, align=PP_ALIGN.LEFT)

    # タイトル下の細いゴールドライン
    SEP_T = TITLE_T + TITLE_H + 60_000
    add_rect(slide, 360_000, SEP_T, W - 720_000, 28_000, C_GOLD)

    # ── 2カラム: サマリー（左）+ 分析（右） ──────────────────
    COL_T  = SEP_T + 120_000
    COL_H  = H - COL_T - 520_000    # フッター分を除く
    GAP    = 120_000
    MARGIN = 360_000
    COL_W  = (W - MARGIN * 2 - GAP) // 2

    # 左列背景
    add_rect(slide, MARGIN, COL_T, COL_W, COL_H, C_LEFT_BG)
    # 左列: ラベルバー（ネイビー）
    LBL_H = 380_000
    add_rect(slide, MARGIN, COL_T, COL_W, LBL_H, C_NAVY_MID)
    add_textbox(slide, MARGIN + 180_000, COL_T + 60_000,
                COL_W - 200_000, LBL_H - 80_000,
                "● 売上サマリー", 11, bold=True, color=C_WHITE)
    # 左列: テキストプレースホルダー
    add_textbox(slide, MARGIN + 120_000, COL_T + LBL_H + 120_000,
                COL_W - 240_000, COL_H - LBL_H - 180_000,
                "{{summary_text}}", 10.5, color=C_TEXT)

    # 右列背景
    R_L = MARGIN + COL_W + GAP
    add_rect(slide, R_L, COL_T, COL_W, COL_H, C_RIGHT_BG)
    # 右列: ラベルバー（ゴールド）
    add_rect(slide, R_L, COL_T, COL_W, LBL_H, C_GOLD)
    add_textbox(slide, R_L + 180_000, COL_T + 60_000,
                COL_W - 200_000, LBL_H - 80_000,
                "● 課題・所見と改善策", 11, bold=True, color=C_NAVY)
    # 右列: テキストプレースホルダー
    add_textbox(slide, R_L + 120_000, COL_T + LBL_H + 120_000,
                COL_W - 240_000, COL_H - LBL_H - 180_000,
                "{{analysis_text}}", 10.5, color=C_TEXT)

    # ── フッター ──────────────────────────────────────────────
    FTR_T = H - 400_000
    add_rect(slide, 0, FTR_T, W, 400_000, C_NAVY)
    # フッター内のゴールドライン（上辺）
    add_rect(slide, 0, FTR_T, W, 28_000, C_GOLD)
    add_textbox(slide, 360_000, FTR_T + 80_000, W // 2, 300_000,
                "© 売上報告書 自動生成システム  |  Powered by Local LLM (Ollama)",
                8, color=RGBColor(0xAA, 0xB8, 0xCC), align=PP_ALIGN.LEFT)
    add_textbox(slide, W // 2, FTR_T + 80_000, W // 2 - 360_000, 300_000,
                "社外秘 — 取扱注意", 8,
                color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.RIGHT)

    out = DATA_DIR / "template_consultant.pptx"
    prs.save(out)
    print(f"テンプレート生成完了: {out}")


# ────────────────────────────────────────────────────────────
#  多年度サンプル売上データ生成
# ────────────────────────────────────────────────────────────
PRODUCTS = [
    # (商品名, 単価, 原価率, 月次出荷傾向)
    ("プレミアムプラン",    150_000, 0.35, [0.8,0.8,1.0,1.0,1.1,1.0,1.0,1.1,1.2,1.2,1.3,1.5]),
    ("スタンダードプラン",   60_000, 0.50, [1.0,0.9,1.0,1.0,1.0,1.1,1.1,1.0,1.0,1.1,1.2,1.4]),
    ("エントリーパッケージ", 25_000, 0.65, [1.1,1.0,1.0,1.0,1.0,1.0,1.0,0.9,1.0,1.0,1.1,1.2]),
    ("コンサルティング",    200_000, 0.30, [0.7,0.8,1.0,1.1,1.1,1.0,0.9,1.0,1.2,1.2,1.1,1.3]),
    ("保守サポート",         18_000, 0.25, [1.0,1.0,1.0,1.0,1.0,1.0,1.0,1.0,1.0,1.0,1.0,1.1]),
]

REGIONS = ["東京", "大阪", "名古屋", "福岡", "札幌"]
REGION_WEIGHT = [0.38, 0.25, 0.17, 0.12, 0.08]

REPS = [
    ("田中 一郎", "東京"),
    ("佐藤 花子", "大阪"),
    ("鈴木 太郎", "東京"),
    ("伊藤 美咲", "名古屋"),
    ("渡辺 健二", "福岡"),
    ("中村 さくら", "東京"),
]

YOY_GROWTH = {2022: 1.00, 2023: 1.12, 2024: 1.19}   # 年間成長率


def gen_rows(year: int, month: int) -> list[dict]:
    rows = []
    count = random.randint(18, 28)
    growth = YOY_GROWTH[year]

    for _ in range(count):
        # 商品選択（プレミアム/コンサルは低頻度）
        prod = random.choices(
            PRODUCTS,
            weights=[0.18, 0.32, 0.28, 0.10, 0.12]
        )[0]
        name, base_price, cost_rate, seasonality = prod

        # 地域・担当者
        region = random.choices(REGIONS, weights=REGION_WEIGHT)[0]
        candidates = [r for r in REPS if r[1] == region]
        if not candidates:
            candidates = REPS
        rep = random.choice(candidates)[0]

        # 数量
        qty = random.randint(1, 5) if base_price >= 100_000 else random.randint(1, 20)

        # 売上金額（季節性 + 年成長 + ランダム揺らぎ）
        seasonal = seasonality[month - 1]
        noise    = random.uniform(0.88, 1.12)
        amount   = int(base_price * qty * seasonal * growth * noise)

        # 原価・利益
        cost   = int(amount * cost_rate * random.uniform(0.95, 1.05))
        profit = amount - cost
        margin = round(profit / amount * 100, 1) if amount > 0 else 0

        # 日付（月内のランダム営業日）
        day = random.randint(1, 28)
        dt  = date(year, month, day)

        rows.append({
            "日付":   dt.strftime("%Y-%m-%d"),
            "商品名": name,
            "担当者": rep,
            "地域":   region,
            "数量":   qty,
            "売上金額": amount,
            "原価":   cost,
            "利益額": profit,
            "利益率(%)": margin,
        })
    return rows


def create_excel():
    all_rows = []
    for year in [2022, 2023, 2024]:
        for month in range(1, 13):
            all_rows.extend(gen_rows(year, month))

    df = pd.DataFrame(all_rows).sort_values("日付").reset_index(drop=True)

    out_xlsx = DATA_DIR / "sample_advanced.xlsx"
    out_csv  = DATA_DIR / "sample_advanced.csv"

    df.to_excel(out_xlsx, index=False)
    df.to_csv(out_csv,   index=False, encoding="utf-8-sig")

    total = df["売上金額"].sum()
    print(f"売上データ生成完了: {len(df)} 行 / 総売上 {total:,.0f}円")
    print(f"  Excel: {out_xlsx}")
    print(f"  CSV  : {out_csv}")


# ────────────────────────────────────────────────────────────
if __name__ == "__main__":
    create_excel()
    create_template()
    print("\nすべてのサンプルファイルを生成しました。")
