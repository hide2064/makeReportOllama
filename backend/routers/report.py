"""
routers/report.py  — ジョブキュー対応版

POST /api/generate            — ジョブをキューに登録し即返却。
GET  /api/progress/{job_id}   — ジョブ状態をポーリングで取得。
GET  /api/download/{job_id}   — 生成済み PPTX をダウンロード。
GET  /api/slides/{job_id}     — スライド内容をプレビュー用 JSON で返す。
GET  /api/queue               — キュー状態の概要を返す。
"""

import logging
import os
import tempfile
import threading
import uuid
from collections import deque
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse

from config import ANALYST_MAX_RETRIES, DATA_DIR, MODEL_ANALYST, MODEL_WRITER, OUTPUT_DIR
from services.excel_reader import read_and_summarize
from services.history_store import append_history
from services.ollama_client import (
    build_analyst_prompt,
    build_writer_prompt,
    check_ollama,
    generate,
    parse_analyst_json,
    parse_writer_response,
)
from services.pptx_generator import generate_pptx
from services.rag_store import search_context

logger = logging.getLogger(__name__)
router = APIRouter()

OUTPUT_DIR.mkdir(exist_ok=True)

# ── ジョブ管理 ────────────────────────────────────────────────
_jobs_lock = threading.Lock()
_jobs: dict[str, dict] = {}       # job_id → ステータス dict
_pending: deque[str]   = deque()  # 実行待ち job_id (FIFO)
MAX_QUEUED = 10   # 受け付けるキュー最大数
MAX_JOBS   = 100  # メモリ上に保持する最大ジョブ数

_executor = ThreadPoolExecutor(max_workers=1)


def _set_job(job_id: str, **kwargs) -> None:
    with _jobs_lock:
        if job_id in _jobs:
            _jobs[job_id].update(kwargs)
    logger.info(f"Job {job_id[:6]}: {kwargs}")


def _cleanup_old_jobs() -> None:
    """完了 / エラーの古いジョブを MAX_JOBS 件まで削減する。"""
    with _jobs_lock:
        finished = [jid for jid, j in _jobs.items() if j.get("done") or j.get("error")]
        for jid in finished[: max(0, len(finished) - MAX_JOBS)]:
            _jobs.pop(jid, None)


# ── バックグラウンド処理 ─────────────────────────────────────
def _run_generation(
    job_id: str,
    excel_data: bytes,
    template_data: Optional[bytes],
    excel_filename: str,
    template_name: str = "",
    slide_options: Optional[dict] = None,
    analyst_model: str = "",
    writer_model: str = "",
    date_from: str = "",
    date_to: str = "",
    extra_context: str = "",
):
    # 待機キューから除去（実行開始）
    with _jobs_lock:
        try:
            _pending.remove(job_id)
        except ValueError:
            pass

    tmpdir = tempfile.mkdtemp()
    try:
        suffix       = Path(excel_filename).suffix.lower() or ".xlsx"
        excel_path   = os.path.join(tmpdir, f"sales_data{suffix}")
        output_path  = str(OUTPUT_DIR / f"report_{job_id}.pptx")
        used_analyst = analyst_model or MODEL_ANALYST
        used_writer  = writer_model  or MODEL_WRITER

        with open(excel_path, "wb") as f:
            f.write(excel_data)

        # テンプレートパスの解決
        if template_name:
            template_path = str(DATA_DIR / template_name)
            if not Path(template_path).exists():
                _set_job(job_id, error=f"テンプレートが見つかりません: {template_name}")
                return
        elif template_data:
            template_path = os.path.join(tmpdir, "template.pptx")
            with open(template_path, "wb") as f:
                f.write(template_data)
        else:
            _set_job(job_id, error="テンプレートが指定されていません。")
            return

        # Step 0: Ollama プリフライトチェック
        _set_job(job_id, step="[0/3]  Ollama の接続とモデルを確認しています...", progress=5)
        try:
            check_ollama(used_analyst)
            check_ollama(used_writer)
        except RuntimeError as e:
            _set_job(job_id, error=str(e))
            return

        # Step 1: Excel 読み込み
        _set_job(job_id, step="[1/3]  Excel / CSV を読み込んでいます...", progress=15)
        try:
            summary_data = read_and_summarize(excel_path, date_from=date_from, date_to=date_to)
        except (FileNotFoundError, ValueError) as e:
            _set_job(job_id, error=str(e))
            return

        # extra_context の受信確認ログ
        if extra_context.strip():
            logger.info(f"extra_context 受信 ({len(extra_context)}字): {extra_context[:200]}")
        else:
            logger.info("extra_context: なし（空）")

        # Step 2a: Analyst AI
        _set_job(job_id,
                 step=f"[2/3]  売上データを解析中です... (Analyst: {used_analyst})\n"
                      "       数値・トレンドを抽出しています。",
                 progress=25)

        def _on_analyst_token(count: int):
            pct = min(25 + count // 5, 45)
            _set_job(job_id,
                     step=f"[2/3]  売上データを解析中です... (Analyst: {used_analyst})\n"
                          f"       生成中... {count} トークン生成済み",
                     progress=pct)

        analyst_data: dict = {}
        analyst_prompt = build_analyst_prompt(summary_data["raw_summary"])
        for attempt in range(1, ANALYST_MAX_RETRIES + 1):
            try:
                analyst_raw  = generate(analyst_prompt, model=used_analyst,
                                        on_token=_on_analyst_token)
                analyst_data = parse_analyst_json(analyst_raw)
            except RuntimeError as e:
                _set_job(job_id, error=str(e))
                return
            if analyst_data:
                logger.info(f"Analyst 完了 (試行 {attempt}): {list(analyst_data.keys())}")
                break
            if attempt < ANALYST_MAX_RETRIES:
                logger.warning(f"Analyst JSON が空 (試行 {attempt})。リトライします。")
                _set_job(job_id,
                         step=f"[2/3]  売上データを再解析中です... (試行 {attempt + 1}/{ANALYST_MAX_RETRIES})",
                         progress=30)

        # Step 2b: RAG
        _set_job(job_id, step="[2/3]  過去レポートから関連情報を検索中です... (RAG)", progress=50)
        rag_context = search_context(summary_data["raw_summary"])
        if rag_context:
            logger.info(f"RAG コンテキスト取得: {len(rag_context)} 字")

        # Step 2c: Writer AI
        _set_job(job_id,
                 step=f"[2/3]  レポート文章を生成中です... (Writer: {used_writer})\n"
                      "       CPU 推論のため数分かかります。このまましばらくお待ちください。",
                 progress=55)

        def _on_writer_token(count: int):
            pct = min(55 + count // 8, 85)
            _set_job(job_id,
                     step=f"[2/3]  レポート文章を生成中です... (Writer: {used_writer})\n"
                          f"       生成中... {count} トークン生成済み",
                     progress=pct)

        try:
            writer_raw = generate(
                build_writer_prompt(
                    analyst_data, summary_data["raw_summary"],
                    rag_context, extra_context,
                ),
                model=used_writer,
                on_token=_on_writer_token,
            )
            summary_text, analysis_text = parse_writer_response(writer_raw)
        except RuntimeError as e:
            _set_job(job_id, error=str(e))
            return

        # Step 3: PPTX 生成
        _set_job(job_id, step="[3/3]  PowerPoint レポートを生成しています...", progress=90)
        try:
            generate_pptx(
                template_path=template_path,
                output_path=output_path,
                summary_text=summary_text,
                analysis_text=analysis_text,
                period=summary_data["period"],
                monthly_totals=summary_data.get("monthly_totals"),
                product_totals=summary_data.get("product_totals"),
                quarterly_product_pivot=summary_data.get("quarterly_product_pivot"),
                quarterly_region_pivot=summary_data.get("quarterly_region_pivot"),
                quarterly_rep_pivot=summary_data.get("quarterly_rep_pivot"),
                monthly_margin=summary_data.get("monthly_margin"),
                slide_options=slide_options,
            )
        except Exception as e:
            _set_job(job_id, error=f"PPTX 生成に失敗しました: {e}")
            return

        append_history(
            job_id=job_id,
            original_filename=excel_filename,
            output_path=output_path,
            analyst_model=used_analyst,
            writer_model=used_writer,
        )
        _set_job(job_id, step="完了しました！", done=True, output_path=output_path, progress=100)
        _cleanup_old_jobs()

    except Exception as e:
        logger.exception("予期しないエラー")
        _set_job(job_id, error=str(e))
    finally:
        import shutil
        shutil.rmtree(tmpdir, ignore_errors=True)


# ── ヘルパー: output_path の解決 ─────────────────────────────
def _resolve_output_path(job_id: str) -> str:
    """_jobs → history_store の順に output_path を解決する。"""
    with _jobs_lock:
        job = _jobs.get(job_id)
    path = job.get("output_path", "") if job else ""
    if path:
        return path
    from services.history_store import get_history_item
    entry = get_history_item(job_id)
    return entry.get("output_path", "") if entry else ""


# ── エンドポイント ─────────────────────────────────────────────

@router.post("/api/generate")
async def generate_report(
    excel_file:          UploadFile           = File(..., description="売上データ Excel/CSV"),
    template_file:       Optional[UploadFile] = File(None, description="PPTX テンプレート"),
    template_name:       str                  = Form(""),
    slide_product_table: bool                 = Form(True),
    slide_region_table:  bool                 = Form(False),
    slide_rep_table:     bool                 = Form(False),
    slide_chart:         bool                 = Form(True),
    chart_product_type:  str                  = Form("bar"),
    analyst_model:       str                  = Form(""),
    writer_model:        str                  = Form(""),
    date_from:           str                  = Form(""),
    date_to:             str                  = Form(""),
    extra_context:       str                  = Form("", description="追加プロンプト（オプション）"),
):
    """ジョブをキューに登録し、即座に job_id を返す。"""
    with _jobs_lock:
        if len(_pending) >= MAX_QUEUED:
            raise HTTPException(
                status_code=429,
                detail=f"キューが満杯です（最大 {MAX_QUEUED} 件）。しばらくお待ちください。",
            )

    job_id         = uuid.uuid4().hex[:12]
    excel_data     = await excel_file.read()
    template_data  = await template_file.read() if template_file else None
    excel_filename = excel_file.filename or "sales_data.xlsx"

    slide_options = {
        "product_table":      slide_product_table,
        "region_table":       slide_region_table,
        "rep_table":          slide_rep_table,
        "chart":              slide_chart,
        "chart_product_type": chart_product_type,
    }

    with _jobs_lock:
        _jobs[job_id] = {
            "step":              "待機中... キューに追加されました",
            "done":              False,
            "error":             "",
            "output_path":       "",
            "progress":          0,
            "original_filename": excel_filename,
        }
        _pending.append(job_id)
        queue_position = len(_pending)

    _executor.submit(
        _run_generation,
        job_id, excel_data, template_data, excel_filename,
        template_name, slide_options, analyst_model, writer_model,
        date_from, date_to, extra_context,
    )
    logger.info(f"ジョブ登録: job_id={job_id}, queue_position={queue_position}")

    return {"status": "queued", "job_id": job_id, "queue_position": queue_position}


@router.get("/api/progress/{job_id}")
async def get_progress(job_id: str):
    """指定ジョブの状態を返す。"""
    with _jobs_lock:
        job = _jobs.get(job_id)
        if not job:
            raise HTTPException(status_code=404, detail="ジョブが見つかりません。")
        try:
            queue_position = list(_pending).index(job_id) + 1
        except ValueError:
            queue_position = 0
        return {**job, "queue_position": queue_position}


@router.get("/api/queue")
async def get_queue():
    """キュー全体の状態を返す。"""
    with _jobs_lock:
        pending_ids = list(_pending)
        active = next(
            (jid for jid, j in _jobs.items()
             if j.get("step") and not j.get("done") and not j.get("error") and jid not in _pending),
            None,
        )
    return {"pending_count": len(pending_ids), "active_job_id": active, "pending_jobs": pending_ids}


@router.get("/api/download/{job_id}")
async def download_report(job_id: str):
    """指定ジョブの PPTX をダウンロードさせる。"""
    output_path = _resolve_output_path(job_id)
    if not output_path or not Path(output_path).exists():
        raise HTTPException(status_code=404, detail="レポートファイルが見つかりません。先に生成してください。")

    with _jobs_lock:
        job = _jobs.get(job_id, {})
    original = job.get("original_filename", "") or ""
    base = Path(original).stem or job_id
    filename = f"report_{base}.pptx"

    return FileResponse(
        path=output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )


@router.get("/api/slides/{job_id}")
async def get_slide_preview(job_id: str):
    """生成済み PPTX のスライド内容を JSON で返す（ブラウザ内プレビュー用）。"""
    output_path = _resolve_output_path(job_id)
    if not output_path or not Path(output_path).exists():
        raise HTTPException(status_code=404, detail="レポートファイルが見つかりません。")

    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation(output_path)
        slides_data = []

        for i, slide in enumerate(prs.slides):
            # タイトルを取得
            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()

            # スライドの種類を判定
            has_table    = any(shape.has_table for shape in slide.shapes)
            has_picture  = any(
                getattr(shape, "shape_type", None) == 13  # MSO_SHAPE_TYPE.PICTURE
                for shape in slide.shapes
            )
            slide_type = "table" if has_table else ("chart" if has_picture else "text")

            # テキスト抽出（タイトル以外）
            texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if text and text != title_text:
                    texts.append(text[:400])

            slides_data.append({
                "slide_num": i + 1,
                "title":     title_text or f"スライド {i + 1}",
                "texts":     texts[:5],
                "type":      slide_type,
            })

        return {"job_id": job_id, "slides": slides_data}

    except Exception as e:
        logger.error(f"スライドプレビュー取得エラー: {e}")
        raise HTTPException(status_code=500, detail=f"スライドの解析に失敗しました: {e}")
