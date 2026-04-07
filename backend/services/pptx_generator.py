"""
pptx_generator.py
テンプレート PPTX のプレースホルダーにテキストを埋め込んで出力する。

Phase 4 追加:
  - 商品別売上表スライド (4a)
  - 月次推移・商品構成グラフスライド (4b)
"""

import io
import logging
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# matplotlib は任意依存: なければグラフスライドをスキップ
try:
    import matplotlib
    matplotlib.use("Agg")
    matplotlib.rcParams.update({
        "font.family":     "sans-serif",
        "font.sans-serif": ["Meiryo", "MS Gothic", "Yu Gothic", "DejaVu Sans"],
    })
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mtick
    _MPL_OK = True
except ImportError:
    _MPL_OK = False

logger = logging.getLogger(__name__)

# ── カラーパレット (エグゼクティブデザイン) ──────────────────────
NAVY       = RGBColor(0x0A, 0x12, 0x28)   # C_INK: deepest navy (cover bg)
NAVY2      = RGBColor(0x16, 0x32, 0x60)   # C_NAVY_MID: section headers
GOLD       = RGBColor(0xD4, 0x94, 0x1A)   # C_GOLD: accent / highlights
TEAL       = RGBColor(0x00, 0xA3, 0x9A)   # C_TEAL: secondary accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xF6, 0xF8, 0xFC)   # C_OFFWHITE: slide backgrounds
TEXT_DARK  = RGBColor(0x0A, 0x12, 0x28)   # same as NAVY for body text
GRAY       = RGBColor(0x55, 0x65, 0x7A)

# ── スライドサイズ (16:9 ワイド) ──────────────────────────────────
W = 12_192_000   # EMU
H =  6_858_000   # EMU


# ── ユーティリティ ────────────────────────────────────────────────

def _get_blank_layout(prs):
    """プレースホルダーなしのレイアウトを探す。なければ最後のレイアウトを返す。"""
    for layout in prs.slide_layouts:
        if not layout.placeholders:
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _add_blank_slide(prs):
    """プレースホルダーを除去した白紙スライドを追加して返す。"""
    layout = _get_blank_layout(prs)
    slide  = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _tb(slide, l, t, w, h, text, size, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Arial"
    return box


def _slide_header(slide, title):
    _rect(slide, 0, 0, W, 680_000, NAVY)
    _rect(slide, 0, 0, 160_000, 680_000, GOLD)
    _tb(slide, 220_000, 100_000, W - 500_000, 480_000,
        title, 18, bold=True, color=WHITE)


def _slide_footer(slide):
    _rect(slide, 0, H - 360_000, W, 360_000, NAVY)
    _rect(slide, 0, H - 360_000, W, 26_000, GOLD)
    _tb(slide, 360_000, H - 310_000, W - 720_000, 280_000,
        "社外秘 — 取扱注意", 8, color=RGBColor(0xAA, 0xB8, 0xCC))


def _set_cell(cell, text, bold=False, bg=None, fg=TEXT_DARK,
              align=PP_ALIGN.LEFT, font_size=9):
    """テーブルセルにテキストとスタイルを設定する。"""
    from pptx.oxml.ns import qn
    tf   = cell.text_frame
    para = tf.paragraphs[0]
    para.alignment = align
    # 既存 run をすべて除去してから追加
    for r in list(para._p.findall(qn("a:r"))):
        para._p.remove(r)
    run = para.add_run()
    run.text           = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = fg
    run.font.name      = "Arial"
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _fmt_man(value: int) -> str:
    """金額を万円単位でフォーマット。例: 3847000 → '385万'"""
    man = value // 10_000
    return f"{man:,}万"


# ── Phase 4a: 売上表スライド ──────────────────────────────────────

def _add_table_slide(prs, pivot, title: str, row_header: str = "商品名"):
    """汎用四半期売上表スライドをプレゼンに追加する。product/region/rep で再利用可能。"""
    if pivot is None or pivot.empty:
        logger.info(f"table_data なし → {title} スライドをスキップ")
        return

    slide = _add_blank_slide(prs)
    _slide_header(slide, title)
    _slide_footer(slide)

    qp       = pivot
    products = list(qp.index)
    quarters = list(qp.columns)

    n_rows = len(products) + 2   # ヘッダー行 + データ行 + 合計行
    n_cols = len(quarters)  + 2  # 商品名列 + 四半期列 + 合計列

    margin      = 360_000
    table_left  = margin
    table_top   = 820_000
    table_width = W - 2 * margin
    table_height = H - table_top - 420_000

    shape = slide.shapes.add_table(
        n_rows, n_cols,
        Emu(table_left), Emu(table_top),
        Emu(table_width), Emu(table_height),
    )
    tbl = shape.table

    # ── 列幅を設定 ────────────────────────────────────────────
    name_w  = 2_200_000
    total_w = 1_400_000
    q_w     = (table_width - name_w - total_w) // max(len(quarters), 1)

    tbl.columns[0].width = Emu(name_w)
    for i in range(len(quarters)):
        tbl.columns[i + 1].width = Emu(q_w)
    tbl.columns[n_cols - 1].width = Emu(total_w)

    # ── 行高を設定 ────────────────────────────────────────────
    tbl.rows[0].height = Emu(520_000)   # ヘッダー行
    for i in range(1, n_rows - 1):
        tbl.rows[i].height = Emu(450_000)
    tbl.rows[n_rows - 1].height = Emu(520_000)  # 合計行

    # ── ヘッダー行 ───────────────────────────────────────────
    _set_cell(tbl.cell(0, 0), row_header, bold=True, bg=NAVY, fg=WHITE,
              align=PP_ALIGN.LEFT, font_size=9)
    for ci, q in enumerate(quarters):
        label = q[2:] if len(q) >= 6 else q   # "2022Q1" → "22Q1"
        _set_cell(tbl.cell(0, ci + 1), label, bold=True, bg=NAVY, fg=WHITE,
                  align=PP_ALIGN.CENTER, font_size=9)
    _set_cell(tbl.cell(0, n_cols - 1), "合計", bold=True, bg=GOLD, fg=WHITE,
              align=PP_ALIGN.CENTER, font_size=9)

    # ── データ行 ─────────────────────────────────────────────
    for ri, product in enumerate(products):
        row_bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
        _set_cell(tbl.cell(ri + 1, 0), product, bg=row_bg, fg=TEXT_DARK,
                  align=PP_ALIGN.LEFT, font_size=9)
        row_total = 0
        for ci, q in enumerate(quarters):
            val = int(qp.loc[product, q])
            row_total += val
            _set_cell(tbl.cell(ri + 1, ci + 1), _fmt_man(val), bg=row_bg,
                      fg=TEXT_DARK, align=PP_ALIGN.RIGHT, font_size=9)
        _set_cell(tbl.cell(ri + 1, n_cols - 1), _fmt_man(row_total),
                  bold=True, bg=row_bg, fg=TEXT_DARK,
                  align=PP_ALIGN.RIGHT, font_size=9)

    # ── 合計行 ───────────────────────────────────────────────
    tr = n_rows - 1
    _set_cell(tbl.cell(tr, 0), "合計", bold=True, bg=NAVY2, fg=GOLD,
              align=PP_ALIGN.LEFT, font_size=9)
    grand_total = 0
    for ci, q in enumerate(quarters):
        col_total = int(qp[q].sum())
        grand_total += col_total
        _set_cell(tbl.cell(tr, ci + 1), _fmt_man(col_total), bold=True,
                  bg=NAVY2, fg=WHITE, align=PP_ALIGN.RIGHT, font_size=9)
    _set_cell(tbl.cell(tr, n_cols - 1), _fmt_man(grand_total), bold=True,
              bg=GOLD, fg=WHITE, align=PP_ALIGN.RIGHT, font_size=9)

    logger.info(f"{title} スライド追加: {len(products)} 行 × {len(quarters)} 四半期")


# ── Phase 4b: グラフスライド ──────────────────────────────────────

def _add_chart_slide(prs, monthly_totals: dict, product_totals: dict,
                     monthly_margin: dict | None = None,
                     product_chart_type: str = "bar"):
    """月次推移バーチャート + 商品別構成横棒グラフをスライドに追加する。"""
    if not _MPL_OK:
        logger.warning("matplotlib が利用不可 → グラフスライドをスキップ")
        return
    if not monthly_totals and not product_totals:
        logger.info("グラフデータなし → グラフスライドをスキップ")
        return

    slide = _add_blank_slide(prs)
    _slide_header(slide, "売上推移グラフ")
    _slide_footer(slide)

    # ── matplotlib 描画（シンプル・コンサル風） ──────────────
    has_margin = bool(monthly_margin and any(
        v == v for v in monthly_margin.values()  # NaN 除外チェック
    ))
    fig, axes = plt.subplots(
        1, 2,
        figsize=(13.0, 4.6),
        gridspec_kw={"width_ratios": [6, 4]},
    )
    fig.patch.set_facecolor("#FFFFFF")

    C_NAVY = "#0A1228"
    C_GOLD = "#D4941A"
    C_TEAL = "#00A39A"
    C_GRAY = "#D1D5DB"

    # 左: 月次売上バーチャート ─────────────────────────────────
    ax1    = axes[0]
    months = list(monthly_totals.keys())
    vals   = [monthly_totals[m] // 10_000 for m in months]
    x_pos  = list(range(len(months)))

    # ネイビー単色バー、エッジなし
    ax1.bar(x_pos, vals, color=C_NAVY, width=0.6, zorder=3)
    ax1.set_xticks(x_pos)
    ax1.set_xticklabels(months, rotation=45, ha="right", fontsize=7)
    ax1.set_ylabel("（万円）", fontsize=8, color="#555")
    ax1.set_title("月次売上推移", fontsize=10, fontweight="bold",
                  color=C_NAVY, pad=10, loc="left")
    ax1.yaxis.set_major_formatter(mtick.FuncFormatter(lambda v, _: f"{v:,.0f}"))
    # グリッドは薄い水平線のみ
    ax1.yaxis.grid(True, linestyle="--", linewidth=0.5, color=C_GRAY, zorder=0)
    ax1.set_axisbelow(True)
    ax1.set_facecolor("#FFFFFF")
    ax1.spines[["top", "right", "left"]].set_visible(False)
    ax1.tick_params(axis="both", length=0)

    # 最高値バーにゴールドマーク
    if vals:
        max_idx = vals.index(max(vals))
        ax1.patches[max_idx].set_facecolor(C_GOLD)
        ax1.text(max_idx, vals[max_idx] + max(vals) * 0.01,
                 f"{vals[max_idx]:,}", ha="center", va="bottom",
                 fontsize=7.5, fontweight="bold", color=C_GOLD)

    # 利益率ライン（オプション）
    if has_margin:
        ax1r = ax1.twinx()
        margin_vals = [monthly_margin.get(m) for m in months]
        ax1r.plot(x_pos, margin_vals, color=C_TEAL, marker="o",
                  markersize=3.5, linewidth=1.5, zorder=4)
        ax1r.set_ylabel("利益率(%)", fontsize=8, color=C_TEAL)
        ax1r.tick_params(axis="y", labelcolor=C_TEAL, length=0, labelsize=7.5)
        ax1r.yaxis.set_major_formatter(mtick.FuncFormatter(lambda v, _: f"{v:.0f}%"))
        ax1r.set_ylim(0, 100)
        ax1r.spines[["top", "right"]].set_visible(False)

    # 右: 商品別売上グラフ ─────────────────────────────────────
    ax2   = axes[1]
    prods = list(product_totals.keys())[:6]   # 上位6件に絞り見やすく
    p_vals = [product_totals[p] // 10_000 for p in prods]

    if product_chart_type == "pie":
        # シンプルドーナツ: ネイビー系グラデーション2色 + ゴールド
        pie_colors = [C_NAVY, "#163260", C_GOLD, C_TEAL, "#556578", "#9CA3AF"][:len(prods)]
        wedges, _, autotexts = ax2.pie(
            p_vals, labels=None, colors=pie_colors[:len(prods)],
            autopct="%1.0f%%", startangle=90, pctdistance=0.80,
            wedgeprops={"width": 0.50, "edgecolor": "white", "linewidth": 2},
        )
        for at in autotexts:
            at.set_fontsize(7.5)
            at.set_color("white")
            at.set_fontweight("bold")
        ax2.legend(wedges, prods, loc="lower center", bbox_to_anchor=(0.5, -0.18),
                   ncol=2, fontsize=7, frameon=False)
        ax2.set_title("商品別売上構成", fontsize=10, fontweight="bold",
                      color=C_NAVY, pad=10, loc="left")
    else:
        # シンプル横棒: ネイビー単色 + 値ラベル
        bar_colors = [C_GOLD if i == 0 else C_NAVY for i in range(len(prods))]
        ax2.barh(range(len(prods)), p_vals, color=bar_colors,
                 height=0.55, zorder=3)
        ax2.set_yticks(range(len(prods)))
        ax2.set_yticklabels(prods, fontsize=8)
        ax2.set_title("商品別売上構成", fontsize=10, fontweight="bold",
                      color=C_NAVY, pad=10, loc="left")
        ax2.xaxis.set_major_formatter(mtick.FuncFormatter(lambda v, _: f"{v:,.0f}"))
        ax2.xaxis.grid(True, linestyle="--", linewidth=0.5, color=C_GRAY, zorder=0)
        ax2.set_axisbelow(True)
        ax2.set_facecolor("#FFFFFF")
        ax2.spines[["top", "right", "left"]].set_visible(False)
        ax2.tick_params(axis="both", length=0)
        ax2.invert_yaxis()
        # 各バーの端に値ラベル
        max_v = max(p_vals) if p_vals else 1
        for bar, val in zip(ax2.patches, p_vals):
            ax2.text(bar.get_width() + max_v * 0.02,
                     bar.get_y() + bar.get_height() / 2,
                     f"{val:,}", va="center", ha="left",
                     fontsize=7.5, color=C_NAVY)

    plt.tight_layout(pad=1.8)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight",
                facecolor="#F6F8FC")
    plt.close(fig)
    buf.seek(0)

    # スライドに画像として埋め込む
    margin      = 180_000
    img_top     = 750_000
    img_width   = W - 2 * margin
    img_height  = H - img_top - 420_000
    slide.shapes.add_picture(buf, Emu(margin), Emu(img_top),
                             Emu(img_width), Emu(img_height))

    logger.info("グラフスライド追加")


# ── Phase 4c: 3年間月次比較グラフスライド ────────────────────────

def _add_multiyear_chart_slide(prs, monthly_by_year: dict):
    """
    3年分の月次売上を月別グループ棒グラフで表示し、
    各年の月平均売上を水平破線で重ねるスライドを追加する。

    monthly_by_year: {year(int): {month_num(int): sales(int), ...}, ...}
    """
    if not _MPL_OK:
        logger.warning("matplotlib が利用不可 → 3年比較グラフをスキップ")
        return
    if not monthly_by_year:
        logger.info("monthly_by_year なし → 3年比較グラフをスキップ")
        return

    slide = _add_blank_slide(prs)
    _slide_header(slide, "月次売上 3年比較")
    _slide_footer(slide)

    years  = sorted(monthly_by_year.keys())
    months = list(range(1, 13))
    month_labels = [f"{m}月" for m in months]

    n_years  = len(years)
    bar_w    = 0.7 / max(n_years, 1)   # グループ内の棒幅
    x_pos    = [m - 1 for m in months]  # 0〜11

    C_NAVY  = "#0A1228"
    C_GOLD  = "#D4941A"
    C_TEAL  = "#00A39A"
    C_GRAY  = "#D1D5DB"
    YEAR_COLORS = [C_NAVY, C_GOLD, C_TEAL]
    LINE_STYLES = ["--", "-.", ":"]

    fig, ax = plt.subplots(figsize=(13.0, 4.6))
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_facecolor("#FFFFFF")

    for yi, year in enumerate(years):
        color    = YEAR_COLORS[yi % len(YEAR_COLORS)]
        offsets  = [x + (yi - n_years / 2 + 0.5) * bar_w for x in x_pos]
        vals_man = [monthly_by_year[year].get(m, 0) // 10_000 for m in months]

        ax.bar(offsets, vals_man, width=bar_w * 0.9, color=color,
               label=str(year), zorder=3, alpha=0.85)

        # 年平均売上（データがある月のみ）
        nonzero = [v for v in vals_man if v > 0]
        if nonzero:
            avg = sum(nonzero) / len(nonzero)
            ax.axhline(avg, color=color, linewidth=1.5,
                       linestyle=LINE_STYLES[yi % len(LINE_STYLES)],
                       label=f"{year}年 月平均 {avg:,.0f}万", zorder=4)

    ax.set_xticks(x_pos)
    ax.set_xticklabels(month_labels, fontsize=8)
    ax.set_ylabel("（万円）", fontsize=8, color="#555")
    ax.set_title("月次売上推移（3年比較）", fontsize=11, fontweight="bold",
                 color=C_NAVY, pad=10, loc="left")
    ax.yaxis.set_major_formatter(
        matplotlib.ticker.FuncFormatter(lambda v, _: f"{v:,.0f}")
    )
    ax.yaxis.grid(True, linestyle="--", linewidth=0.5, color=C_GRAY, zorder=0)
    ax.set_axisbelow(True)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="both", length=0)
    ax.legend(loc="upper left", fontsize=7.5, frameon=False, ncol=n_years * 2)

    plt.tight_layout(pad=1.8)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight",
                facecolor="#F6F8FC")
    plt.close(fig)
    buf.seek(0)

    margin     = 180_000
    img_top    = 750_000
    img_width  = W - 2 * margin
    img_height = H - img_top - 420_000
    slide.shapes.add_picture(buf, Emu(margin), Emu(img_top),
                             Emu(img_width), Emu(img_height))
    logger.info("3年比較グラフスライド追加")


# ── テンプレートプレースホルダー置換 ─────────────────────────────

def _replace_text_in_slide(slide, replacements: dict[str, str]):
    """スライド内のすべてのテキストボックスに対してプレースホルダー置換を行う。"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for placeholder, value in replacements.items():
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, value)
                        logger.debug(f"置換: {placeholder} → (len={len(value)})")


# ── 公開 API ─────────────────────────────────────────────────────

def generate_pptx(
    template_path:             str,
    output_path:               str,
    summary_text:              str,
    analysis_text:             str,
    period:                    str,
    monthly_totals:            dict | None = None,
    product_totals:            dict | None = None,
    quarterly_product_pivot                = None,   # pandas DataFrame | None
    quarterly_region_pivot                 = None,   # pandas DataFrame | None
    quarterly_rep_pivot                    = None,   # pandas DataFrame | None
    monthly_margin:            dict | None = None,
    monthly_by_year:           dict | None = None,
    slide_options:             dict | None = None,
) -> str:
    """
    テンプレートを元に報告書 PPTX を生成して output_path に保存。
    Phase 4: 売上表スライド・グラフスライドをテンプレートの後ろに追加する。
    slide_options キー:
      product_table      (bool, default True)
      region_table       (bool, default False)
      rep_table          (bool, default False)
      chart              (bool, default True)
      chart_product_type ("bar" | "pie", default "bar")
      multiyear_chart    (bool, default False)
    """
    logger.info(f"PPTX 生成開始: template={template_path}")
    if not Path(template_path).exists():
        raise FileNotFoundError(f"テンプレートが見つかりません: {template_path}")

    opts                = slide_options or {}
    do_product_table    = opts.get("product_table", True)
    do_region_table     = opts.get("region_table", False)
    do_rep_table        = opts.get("rep_table", False)
    do_chart            = opts.get("chart", True)
    chart_product_type  = opts.get("chart_product_type", "bar")
    do_multiyear_chart  = opts.get("multiyear_chart", False)

    prs   = Presentation(template_path)
    today = date.today().strftime("%Y年%m月%d日")
    title = f"月次売上報告書（{period}）"

    replacements = {
        "{{report_title}}":  title,
        "{{report_date}}":   f"作成日: {today}",
        "{{summary_text}}":  summary_text,
        "{{analysis_text}}": analysis_text,
    }

    for slide in prs.slides:
        _replace_text_in_slide(slide, replacements)

    # Phase 4a: 売上表スライド群
    if do_product_table:
        _add_table_slide(prs, quarterly_product_pivot, "商品別売上表", "商品名")
    if do_region_table:
        _add_table_slide(prs, quarterly_region_pivot,  "地域別売上表", "地域")
    if do_rep_table:
        _add_table_slide(prs, quarterly_rep_pivot,     "担当者別売上表", "担当者")

    # Phase 4b: グラフスライド
    if do_chart:
        _add_chart_slide(
            prs,
            monthly_totals or {},
            product_totals or {},
            monthly_margin,
            chart_product_type,
        )

    # Phase 4c: 3年比較グラフスライド
    if do_multiyear_chart and monthly_by_year:
        _add_multiyear_chart_slide(prs, monthly_by_year)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info(f"PPTX 保存完了: {output_path}")
    return output_path
