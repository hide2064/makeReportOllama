"""
rag_store.py
ChromaDB を使った過去レポートの保存・検索サービス。

- 過去 PPTX からテキストを抽出してスライド単位でチャンク化
- nomic-embed-text (Ollama) でベクター化して ChromaDB に保存
- 新規レポート生成時に類似チャンクを検索してコンテキストとして返す
"""

import hashlib
import json
import logging
import re
from pathlib import Path

import chromadb
import httpx
from pptx import Presentation

from config import (
    CHROMA_DIR,
    MODEL_EMBED,
    OLLAMA_EMBED_URL,
    OLLAMA_EMBED_TIMEOUT,
    RAG_COLLECTION,
    RAG_MAX_CHUNKS,
    RAG_MAX_CTX_CHARS,
    RAG_MIN_CHUNK,
    RAG_SIM_THRESHOLD,
)

logger = logging.getLogger(__name__)

# ローカルエイリアス（後方互換）
COLLECTION_NAME = RAG_COLLECTION
EMBED_MODEL     = MODEL_EMBED
EMBED_TIMEOUT   = OLLAMA_EMBED_TIMEOUT
MAX_CHUNKS      = RAG_MAX_CHUNKS
MAX_CTX_CHARS   = RAG_MAX_CTX_CHARS
MIN_CHUNK_CHARS = RAG_MIN_CHUNK


# ── ChromaDB クライアント (遅延初期化) ─────────────────────
_client = None
_collection = None


def _get_collection():
    global _client, _collection
    if _collection is None:
        CHROMA_DIR.mkdir(parents=True, exist_ok=True)
        _client = chromadb.PersistentClient(path=str(CHROMA_DIR))
        _collection = _client.get_or_create_collection(
            name=COLLECTION_NAME,
            metadata={"hnsw:space": "cosine"},
        )
    return _collection


# ── 埋め込みベクター生成 ─────────────────────────────────────
def embed_text(text: str) -> list[float]:
    """nomic-embed-text (Ollama) でテキストをベクター化する。"""
    try:
        resp = httpx.post(
            OLLAMA_EMBED_URL,
            json={"model": EMBED_MODEL, "prompt": text},
            timeout=EMBED_TIMEOUT,
        )
        resp.raise_for_status()
        return resp.json()["embedding"]
    except httpx.ConnectError:
        raise RuntimeError(
            "Ollama に接続できません。埋め込みモデルの取得には "
            f"`ollama pull {EMBED_MODEL}` が必要です。"
        )
    except (httpx.HTTPStatusError, KeyError) as e:
        raise RuntimeError(f"埋め込みベクター生成に失敗しました: {e}")


# ── PPTX テキスト抽出 ────────────────────────────────────────
def extract_chunks_from_pptx(pptx_path: str) -> list[str]:
    """
    PPTX からスライド単位でテキストを抽出してチャンクのリストを返す。
    プレースホルダータグ ({{...}}) は除去する。
    """
    prs = Presentation(pptx_path)
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    # プレースホルダーのみの行は除外
                    if line and not re.fullmatch(r"\{\{[^}]+\}\}", line):
                        texts.append(line)
        chunk = "\n".join(texts)
        chunk = re.sub(r"\{\{[^}]+\}\}", "", chunk).strip()
        if len(chunk) >= MIN_CHUNK_CHARS:
            chunks.append(chunk)
            logger.debug(f"  スライド {i}: {len(chunk)} 字")
    return chunks


# ── 過去レポート登録 ─────────────────────────────────────────
def register_report(pptx_path: str, filename: str) -> int:
    """
    PPTX を ChromaDB に登録する。
    同じファイルを再登録した場合は既存データを上書きする。
    登録したチャンク数を返す。
    """
    col = _get_collection()
    chunks = extract_chunks_from_pptx(pptx_path)
    if not chunks:
        logger.warning(f"テキスト抽出できませんでした: {filename}")
        return 0

    # ファイル単位の既存データを削除（上書き対応）
    file_hash = hashlib.md5(filename.encode()).hexdigest()
    try:
        existing = col.get(where={"file_id": file_hash})
        if existing["ids"]:
            col.delete(ids=existing["ids"])
            logger.info(f"既存エントリ削除: {len(existing['ids'])} 件 ({filename})")
    except Exception:
        pass

    ids, embeddings, documents, metadatas = [], [], [], []
    for idx, chunk in enumerate(chunks):
        chunk_id = f"{file_hash}_{idx}"
        try:
            vec = embed_text(chunk)
        except RuntimeError as e:
            logger.error(f"チャンク {idx} のベクター化失敗: {e}")
            raise
        ids.append(chunk_id)
        embeddings.append(vec)
        documents.append(chunk)
        metadatas.append({"filename": filename, "file_id": file_hash, "chunk_idx": idx})

    col.add(ids=ids, embeddings=embeddings, documents=documents, metadatas=metadatas)
    logger.info(f"登録完了: {filename} ({len(ids)} チャンク)")
    return len(ids)


# ── 類似検索 ────────────────────────────────────────────────
def search_context(query: str, n_results: int = MAX_CHUNKS) -> str:
    """
    クエリに近い過去レポートのチャンクを検索し、
    Writer AI に渡す文脈文字列を返す。
    登録件数が 0 の場合は空文字を返す。
    """
    col = _get_collection()
    total = col.count()
    if total == 0:
        logger.info("RAG: 過去資料なし → スキップ")
        return ""

    try:
        query_vec = embed_text(query)
    except RuntimeError as e:
        logger.warning(f"RAG クエリ埋め込み失敗 → スキップ: {e}")
        return ""

    results = col.query(
        query_embeddings=[query_vec],
        n_results=min(n_results, total),
        include=["documents", "metadatas", "distances"],
    )

    chunks_out = []
    total_chars = 0
    for doc, meta, dist in zip(
        results["documents"][0],
        results["metadatas"][0],
        results["distances"][0],
    ):
        similarity = 1 - dist   # cosine distance → similarity
        if similarity < RAG_SIM_THRESHOLD:
            continue
        snippet = doc[:300]     # 1チャンク最大 300 字
        chunks_out.append(
            f"[参考: {meta['filename']} / 類似度 {similarity:.2f}]\n{snippet}"
        )
        total_chars += len(snippet)
        if total_chars >= MAX_CTX_CHARS:
            break

    if not chunks_out:
        return ""

    context = "\n\n".join(chunks_out)
    logger.info(f"RAG: {len(chunks_out)} 件取得 ({total_chars} 字)")
    return context


# ── 登録済みファイル一覧 ─────────────────────────────────────
def list_registered() -> list[dict]:
    """登録済みファイルの一覧を返す。"""
    col = _get_collection()
    if col.count() == 0:
        return []
    all_items = col.get(include=["metadatas"])
    seen: dict[str, dict] = {}
    for meta in all_items["metadatas"]:
        fid = meta["file_id"]
        if fid not in seen:
            seen[fid] = {"filename": meta["filename"], "file_id": fid, "chunks": 0}
        seen[fid]["chunks"] += 1
    return list(seen.values())


# ── チャンク取得 ─────────────────────────────────────────────
def get_chunks_for_file(file_id: str) -> list[dict] | None:
    """指定した file_id のチャンク一覧をスライド順で返す。存在しない場合は None。"""
    col = _get_collection()
    existing = col.get(
        where={"file_id": file_id},
        include=["documents", "metadatas"],
    )
    if not existing["ids"]:
        return None
    chunks = []
    for chunk_id, doc, meta in zip(existing["ids"], existing["documents"], existing["metadatas"]):
        chunks.append({
            "id": chunk_id,
            "text": doc,
            "chunk_idx": meta.get("chunk_idx", 0),
        })
    chunks.sort(key=lambda x: x["chunk_idx"])
    return chunks


# ── ファイル削除 ─────────────────────────────────────────────
def delete_report(file_id: str) -> int:
    """指定した file_id の全チャンクを削除する。削除件数を返す。"""
    col = _get_collection()
    existing = col.get(where={"file_id": file_id})
    if not existing["ids"]:
        return 0
    col.delete(ids=existing["ids"])
    logger.info(f"削除完了: file_id={file_id} ({len(existing['ids'])} チャンク)")
    return len(existing["ids"])
